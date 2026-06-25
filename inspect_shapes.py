from pptx import Presentation

prs = Presentation('gainable_presentation_final.pptx')

with open('shapes_out.txt', 'w', encoding='utf-8') as f:
    for i, slide in enumerate(prs.slides):
        f.write(f"\n=== Slide {i+1} ===\n")
        for shape in slide.shapes:
            f.write(f"Shape Name: {shape.name}, Shape Type: {shape.shape_type}\n")
            if shape.has_text_frame:
                f.write("Text Content:\n")
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        f.write(f"  - {paragraph.text}\n")
            if shape.shape_type == 13: # Picture
                f.write(f"  [Picture] size: {shape.width} x {shape.height}\n")
