import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ==========================================
# COLOR SCHEME
# ==========================================
COLOR_DARK       = RGBColor(31, 45, 61)       # #1F2D3D Navy
COLOR_LIGHT      = RGBColor(248, 250, 252)     # #F8FAFC Slate-50
COLOR_GOLD       = RGBColor(213, 155, 43)      # #D59B2B Gainable Gold
COLOR_WHITE      = RGBColor(255, 255, 255)
COLOR_TEXT_PRI   = RGBColor(31, 45, 61)
COLOR_TEXT_SEC   = RGBColor(100, 116, 139)     # Slate-500
COLOR_LINE       = RGBColor(226, 232, 240)     # Slate-200
COLOR_EMERALD    = RGBColor(16, 185, 129)      # #10B981
COLOR_RED        = RGBColor(239, 68, 68)       # #EF4444
COLOR_ORANGE     = RGBColor(249, 115, 22)      # #F97316
COLOR_BLUE       = RGBColor(59, 130, 246)      # #3B82F6
COLOR_SECTION_BG = RGBColor(241, 245, 249)     # Slate-100

# ==========================================
# ASSET PATHS
# ==========================================
project_root = r"C:\Users\gmaro\.gemini\antigravity-ide\scratch\gainable-fr"
public_dir   = os.path.join(project_root, "public")

path_logo       = os.path.join(public_dir, "logo.png")
path_photo_hvac = os.path.join(public_dir, "hero-hvac.png")
path_photo_diag = os.path.join(public_dir, "diag-inspector.png")
path_photo_vision      = os.path.join(public_dir, "espace_pro_vision_1765140841780.png")
path_photo_conclusion  = os.path.join(public_dir, "espace_pro_conclusion_1765140856043.png")
path_scr_homepage      = os.path.join(project_root, "user_scr_search.png")
path_scr_profile       = os.path.join(project_root, "user_scr_profile.png")
path_scr_articles      = os.path.join(project_root, "user_scr_articles.png")
path_scr_seo           = os.path.join(project_root, "user_scr_seo.png")
path_scr_labels        = os.path.join(project_root, "user_scr_labels.png")
path_scr_pricing       = os.path.join(project_root, "user_scr_pricing.png")
path_chart_seo         = os.path.join(project_root, "seo_impact.png")
path_chart_visibility  = os.path.join(project_root, "seo_curve.png")

# ==========================================
# HELPERS
# ==========================================
def apply_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_separator_line(slide, y_inches, width_start=0.8, width_end=11.733):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(width_start), Inches(y_inches),
        Inches(width_end - width_start), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_LINE
    line.line.color.rgb = COLOR_LINE

def add_header(slide, title_text, category_text="PLAYBOOK COMMERCIAL"):
    # Category tag
    txBox_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.4))
    p_cat = txBox_cat.text_frame.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_GOLD

    # Title
    txBox_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.68), Inches(11.7), Inches(0.9))
    tf = txBox_title.text_frame
    tf.word_wrap = True
    p_title = tf.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Arial"
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_PRI

    add_separator_line(slide, 1.55)

def add_section_badge(slide, x, y, w, h, text, bg_color, text_color=None):
    """Adds a colored badge/pill with text."""
    if text_color is None:
        text_color = COLOR_WHITE
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = bg_color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

def add_script_box(slide, x, y, w, h, title, lines, title_color=None, bg_color=None):
    """Adds a framed script/dialogue box."""
    if title_color is None:
        title_color = COLOR_DARK
    if bg_color is None:
        bg_color = RGBColor(255, 251, 235)  # warm yellow tint

    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg_color
    rect.line.color.rgb = COLOR_GOLD
    rect.line.width = Pt(1.5)

    tx = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.1), Inches(w - 0.3), Inches(h - 0.2))
    tf = tx.text_frame
    tf.word_wrap = True

    # Title
    p_t = tf.paragraphs[0]
    p_t.text = f"💬 {title}"
    p_t.font.name = "Arial"
    p_t.font.size = Pt(9.5)
    p_t.font.bold = True
    p_t.font.color.rgb = title_color
    p_t.space_after = Pt(4)

    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Arial"
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_TEXT_PRI
        p.space_after = Pt(2)

def format_bold_lead(paragraph, bold_text, normal_text, size=Pt(11.5)):
    paragraph.font.name = "Arial"
    paragraph.font.size = size
    run_bold = paragraph.add_run()
    run_bold.text = bold_text
    run_bold.font.bold = True
    run_bold.font.color.rgb = COLOR_TEXT_PRI
    run_normal = paragraph.add_run()
    run_normal.text = normal_text
    run_normal.font.bold = False
    run_normal.font.color.rgb = COLOR_TEXT_SEC

def add_two_col_text(slide, x, y, w, h, items):
    """Add a text box with bullet items."""
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    first = True
    for bold_part, normal_part, spacing in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(spacing)
        format_bold_lead(p, bold_part, normal_part)

# ==========================================
# SLIDE 1: COUVERTURE (Dark Theme)
# ==========================================
slide_layout = prs.slide_layouts[6]  # Blank

slide1 = prs.slides.add_slide(slide_layout)
apply_background(slide1, COLOR_DARK)

# Left panel text
tx_title = slide1.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(6.8), Inches(5.0))
tf_title = tx_title.text_frame
tf_title.word_wrap = True

# Logo
if os.path.exists(path_logo):
    slide1.shapes.add_picture(path_logo, Inches(0.8), Inches(0.8), width=Inches(2.6), height=Inches(0.78))

p_t = tf_title.paragraphs[0]
p_t.text = "Playbook Commercial CVC"
p_t.font.name = "Arial"; p_t.font.size = Pt(38); p_t.font.bold = True
p_t.font.color.rgb = COLOR_WHITE; p_t.space_after = Pt(6)

p_sub = tf_title.add_paragraph()
p_sub.text = "Scripts • Argumentaires • Objections • Closing"
p_sub.font.name = "Arial"; p_sub.font.size = Pt(18); p_sub.font.bold = True
p_sub.font.color.rgb = COLOR_GOLD; p_sub.space_after = Pt(20)

p_desc = tf_title.add_paragraph()
p_desc.text = "Le guide opérationnel complet pour convertir les artisans CVC\nen partenaires Gainable.fr — de la prospection au paiement Stripe."
p_desc.font.name = "Arial"; p_desc.font.size = Pt(13)
p_desc.font.color.rgb = COLOR_WHITE; p_desc.space_after = Pt(20)

# Tags
tags = [
    ("PROSPECTION TÉLÉPHONIQUE", Inches(0.8), Inches(5.2), COLOR_BLUE),
    ("SOCIAL SELLING", Inches(3.8), Inches(5.2), COLOR_EMERALD),
    ("CLOSING STRIPE", Inches(5.8), Inches(5.2), COLOR_GOLD),
]
for tag_text, tx, ty, tc in tags:
    add_section_badge(slide1, tx/Inches(1), ty/Inches(1), 2.5, 0.35, tag_text, tc)

# Right side photo
if os.path.exists(path_photo_hvac):
    slide1.shapes.add_picture(path_photo_hvac, Inches(7.333), Inches(0), width=Inches(6.0), height=Inches(7.5))

# ==========================================
# SLIDE 2: SOMMAIRE
# ==========================================
slide2 = prs.slides.add_slide(slide_layout)
apply_background(slide2, COLOR_LIGHT)
add_header(slide2, "Sommaire du Playbook Commercial CVC", "PLAN")

sections = [
    ("01", "Comprendre le Prospect Artisan CVC", "Les 4 douleurs réelles qui ouvrent la vente", COLOR_BLUE),
    ("02", "Pré-Qualification & Diagnostic Préalable", "Checklist avant tout appel à froid", COLOR_EMERALD),
    ("03", "Profiling Psychologique des Artisans", "Adapter sa posture aux 3 profils types", COLOR_GOLD),
    ("04", "Scripts de Prospection Téléphonique", "Accroche, contournement secrétariat, pitch complet", COLOR_RED),
    ("05", "Social Selling : Instagram & LinkedIn", "Scripts DM Instagram + approches LinkedIn trigger", COLOR_BLUE),
    ("06", "Cadence de Prospection sur 5 Jours", "La séquence mixte hebdomadaire avec 50 prospects", COLOR_EMERALD),
    ("07", "Argumentaires Clés de Vente", "Écoute active, reformulation, anti-plateformes", COLOR_GOLD),
    ("08", "La Démo d'Impact 5 minutes", "Effet WoW IA + Google en direct + Stripe closing", COLOR_RED),
    ("09", "Traitement des Objections", "Les 8 objections critiques avec réponses mot pour mot", COLOR_BLUE),
    ("10", "Plan d'Attaque Hebdomadaire", "Routine du succès en 5 étapes actionnables", COLOR_EMERALD),
    ("11", "Comparatif Plateformes", "Gainable.fr vs Bilik, BNI, Travaux.com, Selocal, PagesJaunes", COLOR_RED),
    ("12", "L'Écosystème de Visibilité", "SEO + Site + Réseaux Sociaux + Google Maps + IA = Maillage", COLOR_BLUE),
    ("13", "Les 3 Phases du Maillage", "Ancrage → Maillage → Dominance locale", COLOR_EMERALD),
    ("14", "Mindset & État d'Esprit", "La conviction, la résilience, le rituel du guerrier", COLOR_GOLD),
    ("15", "C'est un Marathon, Pas un Sprint", "Vision 12 mois — Ligne du temps SEO", COLOR_RED),
]

cols = [sections[:5], sections[5:]]
col_x = [0.8, 7.0]

for ci, col in enumerate(cols):
    for ri, (num, title, sub, color) in enumerate(col):
        y = 1.7 + ri * 1.05
        x = col_x[ci]

        # Number badge
        badge = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(0.5), Inches(0.45))
        badge.fill.solid(); badge.fill.fore_color.rgb = color
        badge.line.color.rgb = color
        tf_b = badge.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = num; p_b.font.name = "Arial"; p_b.font.size = Pt(11); p_b.font.bold = True
        p_b.font.color.rgb = COLOR_WHITE; p_b.alignment = PP_ALIGN.CENTER

        # Title + sub
        tx = slide2.shapes.add_textbox(Inches(x + 0.62), Inches(y), Inches(5.5), Inches(0.5))
        tf = tx.text_frame; tf.word_wrap = True
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Arial"; p_title.font.size = Pt(11); p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_PRI
        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.name = "Arial"; p_sub.font.size = Pt(9)
        p_sub.font.color.rgb = COLOR_TEXT_SEC

# ==========================================
# SLIDE 3: COMPRENDRE LE PROSPECT
# ==========================================
slide3 = prs.slides.add_slide(slide_layout)
apply_background(slide3, COLOR_LIGHT)
add_header(slide3, "Comprendre son Prospect : Les 4 Douleurs de l'Artisan CVC", "01 · PSYCHOLOGIE PROSPECT")

pain_points = [
    (COLOR_RED, "💸", "Revente multiple des leads",
     "Les plateformes (Travaux.com, Effy…) revendent le même prospect à 3-5 concurrents. C'est une guerre des prix qui détruit les marges. L'artisan bat le terrain et finit par brader.",
     "Angle d'attaque : « Chez nous, le client vous appelle vous et vous seul. Zéro concurrent sur votre zone. »"),
    (COLOR_ORANGE, "⏰", "Manque de temps digital",
     "Un artisan passe ses journées sur le chantier. Il n'a ni le temps ni les compétences pour gérer du SEO, des Google Ads ou du contenu web. Il a besoin d'un partenaire clé en main.",
     "Angle d'attaque : « Notre IA rédige vos articles SEO en 2 minutes. Vous n'avez rien à écrire. »"),
    (COLOR_BLUE, "💰", "Coûts publicitaires exorbitants",
     "Sans expertise, Google Ads coûte entre 40 € et 80 € HT par lead climatisation, qu'il soit signé ou non. Un abonnement fixe est bien plus prévisible et rentable.",
     "Angle d'attaque : « Notre abonnement à 850 €/an, c'est 2,33 €/jour. 1 seul chantier gainable (8-15k€) le rentabilise 10 ans. »"),
    (COLOR_EMERALD, "🏝️", "Isolement commercial (TPE vs. Grands groupes)",
     "Face à Maclem, IZI EDF ou Engie, une TPE seule est invisible sur Google. Elle a besoin de l'autorité d'une plateforme nationale pour s'imposer localement.",
     "Angle d'attaque : « Gainable.fr, c'est +58 000 pages référencées. Vous profitez de notre autorité pour dominer votre commune. »"),
]

col_w = 5.7
for i, (color, icon, title, pain, angle) in enumerate(pain_points):
    row = i // 2
    col = i % 2
    x = 0.45 + col * 6.4
    y = 1.7 + row * 2.75

    # Card background
    card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(col_w), Inches(2.45))
    card.fill.solid(); card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = color; card.line.width = Pt(1.5)

    # Top accent bar
    top = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(col_w), Inches(0.08))
    top.fill.solid(); top.fill.fore_color.rgb = color; top.line.color.rgb = color

    # Content
    tx = slide3.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.12), Inches(col_w - 0.4), Inches(2.3))
    tf = tx.text_frame; tf.word_wrap = True

    p_title = tf.paragraphs[0]
    p_title.text = f"{icon}  {title}"
    p_title.font.name = "Arial"; p_title.font.size = Pt(12); p_title.font.bold = True
    p_title.font.color.rgb = color; p_title.space_after = Pt(5)

    p_pain = tf.add_paragraph()
    p_pain.text = pain
    p_pain.font.name = "Arial"; p_pain.font.size = Pt(9.5)
    p_pain.font.color.rgb = COLOR_TEXT_SEC; p_pain.space_after = Pt(5)

    p_angle = tf.add_paragraph()
    p_angle.text = angle
    p_angle.font.name = "Arial"; p_angle.font.size = Pt(9); p_angle.font.bold = True
    p_angle.font.color.rgb = COLOR_DARK

# ==========================================
# SLIDE 4: PRÉ-QUALIFICATION AVANT L'APPEL
# ==========================================
slide4 = prs.slides.add_slide(slide_layout)
apply_background(slide4, COLOR_LIGHT)
add_header(slide4, "Pré-Qualification : La Checklist Avant Tout Appel à Froid", "02 · DIAGNOSTIC PRÉALABLE")

add_section_badge(slide4, 0.8, 1.62, 3.5, 0.32, "⚡ RÈGLE D'OR : Ne jamais appeler à l'aveugle", COLOR_DARK)

checks = [
    ("✅  Présence web basique",
     "Cherchez le nom de l'entreprise sur Google. Site responsive ? Avis GMB récents ? Un site daté = point d'appui clé pour votre pitch.",
     "Commande rapide : tapez le nom de l'entreprise + ville. Si pas de résultat dans les 3 premiers → cible parfaite."),
    ("✅  Audit SEO de la localité",
     "Recherchez 'climatisation gainable + [ville de l'artisan]'. Gainable.fr apparaît ? Montrez-lui qu'en adhérant, il occupe CETTE place.",
     "Si son site n'apparaît pas page 1 sur sa propre ville → argument choc : il est invisible pour ses clients locaux."),
    ("✅  Diagnostic Google Ads",
     "Y a-t-il un badge 'Sponsorisé' sur ses mots-clés ? Oui → il est conscient du digital mais paye trop cher par lead.",
     "Google Ads climatisation gainable coûte entre 3 € et 12 € par clic. Notre abonnement fixe est bien plus rentable."),
    ("✅  Vérification Google Maps",
     "Sa note GMB ? Moins de 3,5 étoiles ou moins de 10 avis → opportunité de valoriser son sérieux via Gainable.",
     "Un artisan RGE avec 4,8★ et 50 avis mais invisible Google est une cible prioritaire : le produit est là, il manque la visibilité."),
    ("✅  Audit rapide site:domaine.com",
     "Tapez 'site:son-site.com' dans Google. Moins de 10 pages indexées → son site est quasi-invisible pour les moteurs.",
     "Gainable.fr = +58 000 pages indexées. L'effet d'autorité est immédiat dès l'adhésion."),
    ("✅  Détection du profil (APE + Dirigeant)",
     "Code APE 4322A (plomberie-chauffage), 4322B (installation CVC), 4339Z (travaux isolation). Dirigeant = bonne cible.",
     "LinkedIn ou Societe.com pour valider : solo artisan, PME 5-20 salariés ou bureau d'études. Adapter le pitch."),
]

col_items = [checks[:3], checks[3:]]
for ci, col in enumerate(col_items):
    x = 0.45 + ci * 6.4
    for ri, (title, detail, tip) in enumerate(col):
        y = 2.1 + ri * 1.67
        tx = slide4.shapes.add_textbox(Inches(x), Inches(y), Inches(5.9), Inches(1.55))
        tf = tx.text_frame; tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.name = "Arial"; p_t.font.size = Pt(11.5); p_t.font.bold = True
        p_t.font.color.rgb = COLOR_DARK; p_t.space_after = Pt(3)

        p_d = tf.add_paragraph()
        p_d.text = detail
        p_d.font.name = "Arial"; p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = COLOR_TEXT_SEC; p_d.space_after = Pt(3)

        p_tip = tf.add_paragraph()
        p_tip.text = f"→ {tip}"
        p_tip.font.name = "Arial"; p_tip.font.size = Pt(9); p_tip.font.bold = True
        p_tip.font.color.rgb = COLOR_EMERALD

        if ri < len(col) - 1:
            add_separator_line(slide4, y + 1.52, x, x + 5.9)

# ==========================================
# SLIDE 5: PROFILING PSYCHOLOGIQUE
# ==========================================
slide5 = prs.slides.add_slide(slide_layout)
apply_background(slide5, COLOR_LIGHT)
add_header(slide5, "Profiling Psychologique : Les 3 Profils Types de l'Artisan CVC", "03 · PSYCHOLOGIE DE VENTE")

profiles = [
    (COLOR_ORANGE, "🔨", "L'Artisan sur Chantier",
     "Le Technicien Pressé",
     ["Il a du bruit, est fatigué, gère tout seul.", "Délai d'attention : max 45 secondes.", "Il raccroche si c'est trop long."],
     "Ultra-direct + jargon technique",
     [
         "« Je fais court, je sais que vous êtes sur un chantier. »",
         "« Vous posez plutôt du Daikin ou du Mitsubishi pour le gainable ? »",
         "« On a une demande exclusive sur [Ville], dispo 2 min ce soir ? »",
     ]),
    (COLOR_RED, "😤", "L'Artisan Réticent",
     "Le Déçu du Web",
     ["A déjà payé des sites inutiles.", "S'est fait piéger par des engagements 36-48 mois.", "Méfiant par défaut envers tout commercial."],
     "Empathie radicale + transparence totale",
     [
         "« Vous avez 100% raison de vous méfier du web. »",
         "« C'est pour ça qu'on fait un tarif fixe annuel, sans engagement abusif de 3 ans. »",
         "« Gainable.fr, c'est pas une plateforme de leads. C'est une vitrine exclusive. »",
     ]),
    (COLOR_BLUE, "👔", "La PME CVC Structurée",
     "Le Dirigeant au Bureau",
     ["Ne pose plus lui-même, gère des équipes.", "Sensible au ROI, au SEO et à l'autorité.", "Cherche à structurer son acquisition."],
     "Professionnelle + ROI sémantique",
     [
         "« Nous renforçons l'autorité sémantique locale de votre PME. »",
         "« Vos concurrents de [Ville] sont déjà sur la plateforme. »",
         "« Je vous prépare une analyse de votre positionnement vs. vos 3 concurrents ? »",
     ]),
]

for i, (color, icon, name, subtitle, traits, posture, scripts) in enumerate(profiles):
    x = 0.35 + i * 4.3
    y = 1.7

    # Profile card
    card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.0), Inches(5.6))
    card.fill.solid(); card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = color; card.line.width = Pt(2.0)

    # Top colored banner
    banner = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(4.0), Inches(1.1))
    banner.fill.solid(); banner.fill.fore_color.rgb = color; banner.line.color.rgb = color

    # Icon + name
    tx_head = slide5.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.05), Inches(3.7), Inches(1.0))
    tf_h = tx_head.text_frame; tf_h.word_wrap = True
    p_icon = tf_h.paragraphs[0]
    p_icon.text = f"{icon}  {name}"
    p_icon.font.name = "Arial"; p_icon.font.size = Pt(13); p_icon.font.bold = True
    p_icon.font.color.rgb = COLOR_WHITE; p_icon.space_after = Pt(2)
    p_sub2 = tf_h.add_paragraph()
    p_sub2.text = subtitle
    p_sub2.font.name = "Arial"; p_sub2.font.size = Pt(10)
    p_sub2.font.color.rgb = RGBColor(255, 255, 200)

    # Body content
    tx_body = slide5.shapes.add_textbox(Inches(x + 0.15), Inches(y + 1.15), Inches(3.7), Inches(4.4))
    tf_b = tx_body.text_frame; tf_b.word_wrap = True

    p_trait_head = tf_b.paragraphs[0]
    p_trait_head.text = "Caractéristiques :"
    p_trait_head.font.name = "Arial"; p_trait_head.font.size = Pt(9.5); p_trait_head.font.bold = True
    p_trait_head.font.color.rgb = color; p_trait_head.space_after = Pt(3)

    for trait in traits:
        p_t2 = tf_b.add_paragraph()
        p_t2.text = f"• {trait}"
        p_t2.font.name = "Arial"; p_t2.font.size = Pt(9)
        p_t2.font.color.rgb = COLOR_TEXT_SEC; p_t2.space_after = Pt(2)

    p_post = tf_b.add_paragraph()
    p_post.text = f"\nPosture : {posture}"
    p_post.font.name = "Arial"; p_post.font.size = Pt(9.5); p_post.font.bold = True
    p_post.font.color.rgb = color; p_post.space_after = Pt(4)

    p_script_head = tf_b.add_paragraph()
    p_script_head.text = "Scripts clés :"
    p_script_head.font.name = "Arial"; p_script_head.font.size = Pt(9.5); p_script_head.font.bold = True
    p_script_head.font.color.rgb = COLOR_DARK; p_script_head.space_after = Pt(3)

    for scr in scripts:
        p_s = tf_b.add_paragraph()
        p_s.text = f"💬 {scr}"
        p_s.font.name = "Arial"; p_s.font.size = Pt(8.5)
        p_s.font.color.rgb = COLOR_TEXT_PRI; p_s.space_after = Pt(3)

# ==========================================
# SLIDE 6: SCRIPTS DE PROSPECTION TÉLÉPHONIQUE
# ==========================================
slide6 = prs.slides.add_slide(slide_layout)
apply_background(slide6, COLOR_LIGHT)
add_header(slide6, "Scripts de Prospection Téléphonique — Mot pour Mot", "04 · SCRIPTS PHONING")

# Left column: Contournement secrétariat + Accroche directe
left_tx = slide6.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(6.0), Inches(5.6))
tf_l = left_tx.text_frame; tf_l.word_wrap = True

p_h1 = tf_l.paragraphs[0]
p_h1.text = "🚧 Contourner le Barrage Secrétariat"
p_h1.font.name = "Arial"; p_h1.font.size = Pt(12); p_h1.font.bold = True
p_h1.font.color.rgb = COLOR_RED; p_h1.space_after = Pt(8)

scripts_secrétariat = [
    ("❌ À ne JAMAIS dire :", "« Je vous appelle pour un partenariat » ou « C'est au sujet du site internet »"),
    ("✅ Formule qui passe :", "« C'est au sujet d'un gainable » → La secrétaire pense à un client chaud qui veut un devis."),
    ("✅ Variante :", "« C'est par rapport à un dossier d'installation gainable sur [Ville]. »"),
    ("✅ Si on vous demande votre société :", "« Je suis de Gainable.fr, c'est un réseau d'installateurs partenaires. »"),
]
for bold, normal in scripts_secrétariat:
    p = tf_l.add_paragraph()
    p.space_after = Pt(5)
    format_bold_lead(p, bold + "  ", normal, size=Pt(10))

add_separator_line(slide6, 3.8, 0.5, 6.5)

p_h2 = tf_l.add_paragraph()
p_h2.text = "\n📞 Script d'Accroche (Gérant décroche)"
p_h2.font.name = "Arial"; p_h2.font.size = Pt(12); p_h2.font.bold = True
p_h2.font.color.rgb = COLOR_BLUE; p_h2.space_after = Pt(8)

accroche_lines = [
    "« Bonjour [Prénom], c'est [Votre Prénom] de Gainable.fr. »",
    "« Je serai rapide. On gère un réseau d'installateurs CVC partenaires exclusifs. »",
    "« J'ai regardé votre présence sur Google ce matin — vous n'apparaissez pas »",
    "  « sur 'gainable [Ville]' alors que des chantiers à 10-15k€ se signent sur ces recherches. »",
    "« On a une place disponible sur votre secteur. Ça prend 5 min pour vous montrer. »",
    "« Vous êtes dispo quand cette semaine ? »",
]
for line in accroche_lines:
    p = tf_l.add_paragraph()
    p.text = line
    p.font.name = "Arial"; p.font.size = Pt(10)
    if line.startswith("«"):
        p.font.color.rgb = COLOR_DARK; p.font.bold = False
    else:
        p.font.color.rgb = COLOR_TEXT_SEC
    p.space_after = Pt(3)

# Right column: Script long + rebond DM
right_tx = slide6.shapes.add_textbox(Inches(6.9), Inches(1.7), Inches(6.0), Inches(5.6))
tf_r = right_tx.text_frame; tf_r.word_wrap = True

p_h3 = tf_r.paragraphs[0]
p_h3.text = "📋 Script Long — Qualification Complète (3 min)"
p_h3.font.name = "Arial"; p_h3.font.size = Pt(12); p_h3.font.bold = True
p_h3.font.color.rgb = COLOR_EMERALD; p_h3.space_after = Pt(8)

long_script = [
    ("INTRO :", "« Vous installez principalement du gainable ou vous faites aussi du split mural ? »"),
    ("ÉCOUTE :", "Laissez-le parler. Notez les marques (Daikin, Mitsubishi, Fujitsu), les zones, les volumes."),
    ("DOULEUR :", "« Comment vous gérez l'acquisition de nouveaux chantiers actuellement ? »"),
    ("PONT :", "« Vous connaissez le stress de rappeler un client en 3 min avant 4 concurrents ? »"),
    ("PIVOT :", "« Chez nous, c'est l'inverse. Le client voit VOTRE fiche, VOTRE numéro, et il vous appelle directement. »"),
    ("STATS :", "« Sur Gainable.fr, on a +58 000 pages indexées. Votre ville, votre nom, en tête de Google. »"),
    ("CLOSE :", "« Je vous fais une démo en direct de 5 min pour que vous voyez comment ça marche ? »"),
    ("URGENCE :", "« J'ai une dernière place sur [Ville] avant de la proposer à [Concurrent Connu si possible]. »"),
]
for bold, normal in long_script:
    p = tf_r.add_paragraph()
    p.space_after = Pt(5)
    format_bold_lead(p, f"  {bold}  ", normal, size=Pt(10))

add_separator_line(slide6, 6.7, 6.9, 12.9)

p_tip = tf_r.add_paragraph()
p_tip.text = "\n💡 Si le gérant rappelle suite à un DM envoyé la veille :"
p_tip.font.name = "Arial"; p_tip.font.size = Pt(10); p_tip.font.bold = True
p_tip.font.color.rgb = COLOR_ORANGE; p_tip.space_after = Pt(4)
p_tip2 = tf_r.add_paragraph()
p_tip2.text = "« Ah oui ! Je vous avais envoyé un message hier sur [Instagram/LinkedIn]. Merci de me rappeler. Je voulais vous parler de la zone exclusive sur [Ville] avant qu'on la libère à vos confrères. »"
p_tip2.font.name = "Arial"; p_tip2.font.size = Pt(10)
p_tip2.font.color.rgb = COLOR_TEXT_PRI; p_tip2.font.bold = True

# ==========================================
# SLIDE 7: SOCIAL SELLING — INSTAGRAM & LINKEDIN
# ==========================================
slide7 = prs.slides.add_slide(slide_layout)
apply_background(slide7, COLOR_LIGHT)
add_header(slide7, "Social Selling : Scripts Instagram & LinkedIn — Mot pour Mot", "05 · RÉSEAUX SOCIAUX")

# ---- Left: Instagram ----
insta_tx = slide7.shapes.add_textbox(Inches(0.5), Inches(1.65), Inches(6.0), Inches(5.6))
tf_insta = insta_tx.text_frame; tf_insta.word_wrap = True

p_insta_title = tf_insta.paragraphs[0]
p_insta_title.text = "📸 Instagram — Aucun Barrage Secrétaire"
p_insta_title.font.name = "Arial"; p_insta_title.font.size = Pt(12.5); p_insta_title.font.bold = True
p_insta_title.font.color.rgb = RGBColor(225, 48, 108); p_insta_title.space_after = Pt(6)

insta_scripts = [
    ("Étape 1 — Story-Réaction Technique",
     "Réagissez à sa story (pose gainable, chantier, clim en cours). Question de métier :",
     "« Propre la pose ! C'est du bi-split Daikin ? Pas trop galère le passage des liaisons ? »",
     "L'artisan adore parler technique. Ça brise la glace sans approche commerciale."),
    ("Étape 2 — DM Audit Visuel",
     "Après interaction, envoyez en DM :",
     "« Ton boulot est super propre. C'est dommage car quand je cherche 'clim gainable [Ville]', tu n'apparais pas. Les gros chantiers à 12k€ se signent sur Google, pas sur Insta. On a une demande exclusive sur ton secteur. Dispo pour s'appeler 2 min ? »",
     ""),
    ("Étape 3 — Relance si pas de réponse (J+2)",
     "Message court de relance :",
     "« Salut [Prénom], je voulais valider ton intérêt pour la zone [Ville] avant de la libérer pour tes confrères. »",
     "L'urgence est réelle : l'exclusivité par zone est un vrai argument."),
]

for i, (etape, context, script, tip) in enumerate(insta_scripts):
    p_etape = tf_insta.add_paragraph() if i > 0 else tf_insta.add_paragraph()
    if i > 0:
        p_etape.text = ""
        p_etape.space_after = Pt(2)
        p_etape = tf_insta.add_paragraph()

    p_etape.text = etape
    p_etape.font.name = "Arial"; p_etape.font.size = Pt(10.5); p_etape.font.bold = True
    p_etape.font.color.rgb = COLOR_DARK; p_etape.space_after = Pt(3)

    if context:
        p_ctx = tf_insta.add_paragraph()
        p_ctx.text = context
        p_ctx.font.name = "Arial"; p_ctx.font.size = Pt(9.5)
        p_ctx.font.color.rgb = COLOR_TEXT_SEC; p_ctx.space_after = Pt(3)

    p_scr = tf_insta.add_paragraph()
    p_scr.text = f"💬 {script}"
    p_scr.font.name = "Arial"; p_scr.font.size = Pt(9.5); p_scr.font.bold = True
    p_scr.font.color.rgb = COLOR_TEXT_PRI; p_scr.space_after = Pt(3)

    if tip:
        p_tip3 = tf_insta.add_paragraph()
        p_tip3.text = f"→ {tip}"
        p_tip3.font.name = "Arial"; p_tip3.font.size = Pt(9)
        p_tip3.font.color.rgb = COLOR_EMERALD; p_tip3.space_after = Pt(4)

# ---- Right: LinkedIn ----
add_separator_line(slide7, 1.62, 6.6, 13.1)

li_tx = slide7.shapes.add_textbox(Inches(6.9), Inches(1.65), Inches(6.0), Inches(5.6))
tf_li = li_tx.text_frame; tf_li.word_wrap = True

p_li_title = tf_li.paragraphs[0]
p_li_title.text = "💼 LinkedIn — Cibler les PME CVC (5-20 salariés)"
p_li_title.font.name = "Arial"; p_li_title.font.size = Pt(12.5); p_li_title.font.bold = True
p_li_title.font.color.rgb = RGBColor(0, 119, 181); p_li_title.space_after = Pt(6)

li_scripts = [
    ("Approche 1 — Trigger Recrutement",
     "Ciblez les dirigeants qui recrutent sur LinkedIn (ex: 'frigoriste CVC [Ville]') :",
     "« Bonjour [Nom], j'ai vu que vous recrutiez un frigoriste sur [Ville], signe d'une belle croissance ! Dans le cadre de notre réseau gainable.fr, nous cherchons notre installateur partenaire exclusif sur votre secteur. Notre modèle : un unique installateur par zone géographique. Vos avis et certifications témoignent de votre sérieux. Êtes-vous en mesure de prendre de nouveaux chantiers ? »",
     ""),
    ("Approche 2 — Exclusivité Sectorielle",
     "Connexion directe sur profil dirigeant CVC :",
     "« Bonjour [Nom], nous finalisons le référencement exclusif des experts CVC sur [Ville]. Notre modèle : un unique installateur par secteur géographique pour lui envoyer toutes nos demandes sans revente multiple. Êtes-vous en mesure de prendre de nouveaux chantiers ? »",
     ""),
    ("Astuce — Invitation sans note de connexion",
     "Envoyez des demandes de connexion SANS note de message.",
     "Taux d'acceptation supérieur de +30% car ça ne ressemble pas à une sollicitation commerciale directe.",
     "Attendez l'acceptation, puis envoyez le DM personnalisé dans les 24h."),
]

for i, (etape, context, script, tip) in enumerate(li_scripts):
    if i > 0:
        p_blank = tf_li.add_paragraph()
        p_blank.text = ""; p_blank.space_after = Pt(2)

    p_e = tf_li.add_paragraph()
    p_e.text = etape
    p_e.font.name = "Arial"; p_e.font.size = Pt(10.5); p_e.font.bold = True
    p_e.font.color.rgb = COLOR_DARK; p_e.space_after = Pt(3)

    if context:
        p_c = tf_li.add_paragraph()
        p_c.text = context
        p_c.font.name = "Arial"; p_c.font.size = Pt(9.5)
        p_c.font.color.rgb = COLOR_TEXT_SEC; p_c.space_after = Pt(3)

    p_s = tf_li.add_paragraph()
    p_s.text = f"💬 {script}"
    p_s.font.name = "Arial"; p_s.font.size = Pt(9.5); p_s.font.bold = True
    p_s.font.color.rgb = COLOR_TEXT_PRI; p_s.space_after = Pt(3)

    if tip:
        p_t4 = tf_li.add_paragraph()
        p_t4.text = f"→ {tip}"
        p_t4.font.name = "Arial"; p_t4.font.size = Pt(9)
        p_t4.font.color.rgb = RGBColor(0, 119, 181); p_t4.space_after = Pt(4)

# ==========================================
# SLIDE 8: CADENCE DE PROSPECTION 5 JOURS
# ==========================================
slide8 = prs.slides.add_slide(slide_layout)
apply_background(slide8, COLOR_LIGHT)
add_header(slide8, "La Cadence de Prospection : Séquence Mixte sur 5 Jours", "06 · CADENCE & RYTHME")

add_section_badge(slide8, 0.5, 1.62, 5.5, 0.32,
    "🎯 Objectif : Traiter 50 prospects CVC RGE qualifiés par semaine", COLOR_DARK)

days = [
    ("J1", "Mardi", COLOR_BLUE, "Premier Contact Multi-canal", [
        "Likez / commentez une photo technique de pose récente.",
        "Envoyez une demande de connexion (sans note) sur Instagram & LinkedIn.",
        "DM story-réaction : « Propre la pose ! C'est du bi-split Daikin ? »",
        "Si email disponible → Séquence email J1 (présentation + audit offert).",
    ]),
    ("J2", "Mercredi", COLOR_EMERALD, "Premier Appel Téléphonique", [
        "Appel direct entre 8h30-9h30 ou 12h00-13h30.",
        "Barrage secrétariat : « C'est au sujet d'un gainable ».",
        "Rebond DM si déjà contacté : « Je vous avais envoyé un message hier ».",
        "Si messagerie vocale : Laissez votre prénom + gainable.fr + rappel.",
    ]),
    ("J3", "Jeudi", COLOR_ORANGE, "Relance Écrite Courte", [
        "Message court sur LinkedIn / Instagram / WhatsApp :",
        "« Bonjour [Prénom], je vous ai laissé un message hier. »",
        "« Je voulais valider votre intérêt pour l'exclusivité de [Ville] »",
        "« avant de libérer la zone à vos confrères. »",
    ]),
    ("J4", "Vendredi", COLOR_GOLD, "Second Appel / Close Direct", [
        "Deuxième appel — mentionnez le message du J3.",
        "Proposez la démo de 5 min immédiate : « Je suis là, dispo maintenant. »",
        "Si accord → envoyez le lien Stripe directement en ligne avec lui.",
        "Si hésitation → code EXCLU10 pour -10% si signature ce jour.",
    ]),
    ("J5", "Lundi +1", COLOR_RED, "Retrait Stratégique / Urgence Finale", [
        "Message de désengagement = déclencheur puissant :",
        "« Bonjour [Nom], sans retour de votre part, »",
        "« je libère la zone de [Ville] pour vos confrères. »",
        "« Excellente continuation. »",
    ]),
]

col_count = 5
col_w_d = 12.333 / col_count
for i, (num, day, color, title, actions) in enumerate(days):
    x = 0.5 + i * (col_w_d + 0.12)
    y = 2.1

    # Column header
    header_card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(col_w_d), Inches(0.85))
    header_card.fill.solid(); header_card.fill.fore_color.rgb = color
    header_card.line.color.rgb = color

    tx_h = slide8.shapes.add_textbox(Inches(x + 0.05), Inches(y + 0.05), Inches(col_w_d - 0.1), Inches(0.75))
    tf_h2 = tx_h.text_frame
    p_num = tf_h2.paragraphs[0]
    p_num.text = f"{num} · {day}"
    p_num.font.name = "Arial"; p_num.font.size = Pt(11); p_num.font.bold = True
    p_num.font.color.rgb = COLOR_WHITE; p_num.alignment = PP_ALIGN.CENTER; p_num.space_after = Pt(2)
    p_t_d = tf_h2.add_paragraph()
    p_t_d.text = title
    p_t_d.font.name = "Arial"; p_t_d.font.size = Pt(8.5); p_t_d.font.bold = True
    p_t_d.font.color.rgb = RGBColor(255, 255, 200); p_t_d.alignment = PP_ALIGN.CENTER

    # Action cards
    for j, action in enumerate(actions):
        ay = y + 0.95 + j * 1.12
        action_bg = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(ay), Inches(col_w_d), Inches(1.05))
        action_bg.fill.solid(); action_bg.fill.fore_color.rgb = COLOR_WHITE
        action_bg.line.color.rgb = color; action_bg.line.width = Pt(1)

        tx_a = slide8.shapes.add_textbox(
            Inches(x + 0.1), Inches(ay + 0.08),
            Inches(col_w_d - 0.2), Inches(0.9))
        tf_a = tx_a.text_frame; tf_a.word_wrap = True
        p_a = tf_a.paragraphs[0]
        p_a.text = action
        p_a.font.name = "Arial"; p_a.font.size = Pt(8.5)
        p_a.font.color.rgb = COLOR_TEXT_PRI

# ==========================================
# SLIDE 9: ARGUMENTAIRES DE VENTE
# ==========================================
slide9 = prs.slides.add_slide(slide_layout)
apply_background(slide9, COLOR_LIGHT)
add_header(slide9, "Argumentaires Clés : Écoute Active, Reformulation, Anti-Plateformes", "07 · POSTURE & ARGUMENTATION")

args = [
    (COLOR_BLUE, "🎧 1. L'Écoute Active (70% de la vente)", [
        "Valider les frustrations avec empathie :",
        "« Je comprends, c'est très frustrant de payer pour des leads partagés avec 4 concurrents. »",
        "Poser des questions ouvertes :",
        "« Comment gérez-vous les périodes creuses ? » · « Combien vous coûte l'acquisition d'un client ? »",
        "Reformulation pour déclencher le 'Oui franc' :",
        "« Si je comprends bien : vous voulez un canal direct, sans intermédiaire, sans destruction de marges ? »",
    ]),
    (COLOR_EMERALD, "⚖️ 2. Le Comparatif ROI — Chiffres qui Claquent", [
        "Ramener le tarif à l'échelle quotidienne :",
        "850 €/an = 2,33 €/jour = moins qu'un café quotidien.",
        "Comparer avec Google Ads :",
        "Un lead climatisation Google Ads = 40 à 80 € HT (qu'il soit signé ou non).",
        "Le chantier ROI :",
        "1 seule pose gainable (8 000 à 15 000 €) = abonnement rentabilisé pour 10 ans.",
    ]),
    (COLOR_GOLD, "🏆 3. L'Anti-Plateforme — Notre Différence Fondamentale", [
        "Nous ne sommes PAS une plateforme de leads :",
        "Parcours client : Google → Gainable.fr → Votre fiche → Il vous appelle DIRECTEMENT.",
        "0% de commission sur vos chantiers. Abonnement fixe et prévisible.",
        "Argument autorité :",
        "+58 000 pages indexées. Tapez 'site:gainable.fr' → C'est votre bouclier contre les géants.",
        "Fiches d'intervention détaillées → Moins de 'curieux', plus de clients sérieux.",
    ]),
    (COLOR_RED, "🎯 4. Le Closing à la Reformulation Finale", [
        "Après avoir géré les objections, recentrez sur le 'Oui' :",
        "« Donc si je comprends bien, votre seul frein c'est [l'objection nommée] ? »",
        "« Si je règle ça maintenant, on peut avancer ? »",
        "Envoyez le lien Stripe en direct :",
        "Restez à l'oral avec lui pendant la saisie. Si hésitation après clic :",
        "« Je vous applique le code EXCLU10 pour -10% si vous validez maintenant → 765 € au lieu de 850 € HT. »",
    ]),
]

col_pos = [(0.4, 6.1), (6.7, 6.1)]
row_pos = [1.7, 4.15]

for i, (color, title, points) in enumerate(args):
    col = i % 2
    row = i // 2
    x = col_pos[col][0]
    y = row_pos[row]
    w = col_pos[col][1]
    h = 2.3

    card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = color; card.line.width = Pt(1.5)

    top_bar = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.07))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = color; top_bar.line.color.rgb = color

    tx = slide9.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.12), Inches(w - 0.4), Inches(h - 0.22))
    tf = tx.text_frame; tf.word_wrap = True

    p_head = tf.paragraphs[0]
    p_head.text = title
    p_head.font.name = "Arial"; p_head.font.size = Pt(11); p_head.font.bold = True
    p_head.font.color.rgb = color; p_head.space_after = Pt(5)

    for j, pt in enumerate(points):
        p = tf.add_paragraph()
        p.text = pt
        p.font.name = "Arial"
        if pt.startswith("«") or pt.startswith("1 seul") or pt.startswith("850") or pt.startswith("Un lead") or pt.startswith("+58") or pt.startswith("0%") or pt.startswith("Fich"):
            p.font.size = Pt(9.5); p.font.bold = True; p.font.color.rgb = COLOR_TEXT_PRI
        else:
            p.font.size = Pt(9); p.font.bold = False; p.font.color.rgb = COLOR_TEXT_SEC
        p.space_after = Pt(2)

# ==========================================
# SLIDE 10: DÉMO 5 MINUTES + CLOSING STRIPE
# ==========================================
slide10 = prs.slides.add_slide(slide_layout)
apply_background(slide10, COLOR_LIGHT)
add_header(slide10, "La Démo d'Impact 5 min & Le Closing Stripe — Pas à Pas", "08 · DÉMO & CLOSING")

steps = [
    ("⏱ 0:00-2:00", COLOR_BLUE, "Effet WoW de l'IA", [
        "« Dites-moi votre dernier chantier en 3 mots. »",
        "Tapez en direct dans l'assistant IA : 'Pose gainable Daikin à [Ville]'.",
        "L'IA génère un article technique de 500 mots en 30 secondes sous ses yeux.",
        "« Voilà ce que Google verra dans les prochains jours. »",
    ]),
    ("⏱ 2:00-3:00", COLOR_EMERALD, "Preuve : 0 Revente de Contact", [
        "Montrez une fiche artisan Gainable.fr déjà en ligne.",
        "« Dès que le client vous trouve, il clique sur VOTRE numéro — personne d'autre. »",
        "Montrez un exemple de fiche pro déjà en ligne (Air G Énergie par ex.).",
        "« Ces artisans reçoivent des appels directs. Le client les a CHOISIS, pas un algo qui revend. »",
    ]),
    ("⏱ 3:00-4:00", COLOR_GOLD, "Position Google en Direct", [
        "Tapez 'climatisation gainable [Ville cible]' sur Google.",
        "Gainable.fr apparaît en haut. Montrez-le-lui.",
        "« Vous voyez ? Votre profil sera là dans 48h. »",
        "« Vos concurrents qui apparaissent ici ? Ce sont eux que vos clients appellent aujourd'hui à votre place. »",
    ]),
    ("⏱ 4:00-5:00", COLOR_RED, "Closing & Paiement Stripe", [
        "Envoyez le lien de souscription Stripe EN DIRECT.",
        "Restez en ligne avec lui pendant la saisie. Ne raccrochez pas.",
        "Si hésitation : « Je vous applique EXCLU10 → 765 € au lieu de 850 € HT. »",
        "« C'est valable uniquement maintenant, cette zone part dès demain. »",
    ]),
]

for i, (timing, color, title, actions) in enumerate(steps):
    x = 0.4 + i * 3.2
    y = 1.7

    # Step number card
    num_card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.95), Inches(0.6))
    num_card.fill.solid(); num_card.fill.fore_color.rgb = color; num_card.line.color.rgb = color

    tx_n = slide10.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.08), Inches(2.75), Inches(0.5))
    tf_n = tx_n.text_frame
    p_n = tf_n.paragraphs[0]
    p_n.text = f"{timing}"
    p_n.font.name = "Arial"; p_n.font.size = Pt(10); p_n.font.bold = True
    p_n.font.color.rgb = COLOR_WHITE; p_n.alignment = PP_ALIGN.CENTER

    # Title card
    title_card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y + 0.65), Inches(2.95), Inches(0.5))
    title_card.fill.solid(); title_card.fill.fore_color.rgb = RGBColor(241, 245, 249)
    title_card.line.color.rgb = color; title_card.line.width = Pt(1.5)

    tx_title_s = slide10.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.7), Inches(2.75), Inches(0.4))
    tf_ts = tx_title_s.text_frame
    p_ts = tf_ts.paragraphs[0]
    p_ts.text = title
    p_ts.font.name = "Arial"; p_ts.font.size = Pt(11); p_ts.font.bold = True
    p_ts.font.color.rgb = color; p_ts.alignment = PP_ALIGN.CENTER

    # Action items
    for j, action in enumerate(actions):
        ay = y + 1.22 + j * 1.32
        action_card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(ay), Inches(2.95), Inches(1.22))
        action_card.fill.solid(); action_card.fill.fore_color.rgb = COLOR_WHITE
        action_card.line.color.rgb = color; action_card.line.width = Pt(1)

        tx_a2 = slide10.shapes.add_textbox(Inches(x + 0.12), Inches(ay + 0.1), Inches(2.71), Inches(1.05))
        tf_a2 = tx_a2.text_frame; tf_a2.word_wrap = True
        p_a2 = tf_a2.paragraphs[0]
        p_a2.text = f"{j+1}. {action}"
        p_a2.font.name = "Arial"; p_a2.font.size = Pt(9.5)
        if action.startswith("«"):
            p_a2.font.bold = True; p_a2.font.color.rgb = COLOR_TEXT_PRI
        else:
            p_a2.font.color.rgb = COLOR_TEXT_SEC

# ==========================================
# SLIDE 11: TRAITEMENT DES OBJECTIONS
# ==========================================
slide11 = prs.slides.add_slide(slide_layout)
apply_background(slide11, COLOR_LIGHT)
add_header(slide11, "Traitement des Objections : 8 Réponses Mot pour Mot", "09 · OBJECTIONS & CONTRE-ARGUMENTS")

objections = [
    ("💬 « Je fonctionne au bouche-à-oreille »",
     "« C'est génial, preuve de votre sérieux. Mais 9 clients sur 10 recommandés par un ami tapent votre nom sur Google pour se rassurer. S'ils ne voient rien de récent, ils doutent. On digitalise votre bouche-à-oreille. »",
     COLOR_ORANGE),
    ("💬 « J'ai déjà mon propre site vitrine »",
     "« C'est très bien. Tapez 'site:votresite.com' : il n'y a que 4 pages. Face aux géants nationaux, un petit site individuel est invisible. En publiant chez nous, vous profitez de l'autorité de Gainable.fr (58k pages). »",
     COLOR_BLUE),
    ("💬 « J'achète des leads chez Travaux.com/Effy »",
     "« Vous connaissez le stress de rappeler en 3 min avant 4 concurrents et brader vos marges ? Sur Gainable.fr, le client voit VOTRE fiche, il vous choisit, il vous appelle directement. Zéro revente de contact. 0% commission. »",
     COLOR_RED),
    ("💬 « C'est trop cher »",
     "« Combien vous rapporte une seule pose de gainable ? Entre 8k€ et 15k€. Il vous suffit de signer un seul chantier via Gainable.fr pour rentabiliser votre adhésion pour les 10 prochaines années. »",
     COLOR_GOLD),
    ("💬 « Je n'ai pas le temps pour les articles SEO »",
     "« Vous n'avez rien à écrire. Notre IA génère un article de 800 mots en 2 minutes. Vous saisissez le nom de la ville + 3 mots sur le chantier. L'IA fait le reste, optimisé SEO, prêt à publier en 1 clic. »",
     COLOR_EMERALD),
    ("💬 « Je suis complet pour 6 mois »",
     "« Parfait ! Le SEO prend 3-6 mois. Commencer maintenant permet de blinder votre carnet pour la prochaine saison. Et avoir plus de demandes vous permet de choisir les chantiers les plus rentables. »",
     COLOR_BLUE),
    ("💬 « Je vais y réfléchir »",
     "« Je comprends. Pendant que vous réfléchissez, vos clients tapent 'gainable [Ville]' sur Google et appellent un artisan concurrent qui s'est décidé avant vous. Qu'est-ce qui vous retient vraiment ? » (Identifier l'objection cachée)",
     COLOR_RED),
    ("💬 « Je n'ai jamais entendu parler de vous »",
     "« C'est normal, on est récent et on grandit vite. Tapez 'climatisation gainable [Ville]' sur Google maintenant. On est en première page. Vos futurs clients, eux, nous voient déjà. »",
     COLOR_ORANGE),
]

# 4 objections par colonne
col_o = [(0.4, 6.1), (6.7, 6.1)]
for i, (obj_text, response, color) in enumerate(objections):
    col = i % 2
    row = i // 2
    x = col_o[col][0]
    y = 1.7 + row * 1.38

    card = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(col_o[col][1]), Inches(1.28))
    card.fill.solid(); card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = color; card.line.width = Pt(1.5)

    left_bar = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(0.08), Inches(1.28))
    left_bar.fill.solid(); left_bar.fill.fore_color.rgb = color; left_bar.line.color.rgb = color

    tx = slide11.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.08), Inches(col_o[col][1] - 0.28), Inches(1.12))
    tf = tx.text_frame; tf.word_wrap = True

    p_obj = tf.paragraphs[0]
    p_obj.text = obj_text
    p_obj.font.name = "Arial"; p_obj.font.size = Pt(9.5); p_obj.font.bold = True
    p_obj.font.color.rgb = color; p_obj.space_after = Pt(4)

    p_rep = tf.add_paragraph()
    p_rep.text = f"↳ {response}"
    p_rep.font.name = "Arial"; p_rep.font.size = Pt(9)
    p_rep.font.color.rgb = COLOR_TEXT_PRI

# ==========================================
# SLIDE 12: PLAN D'ATTAQUE HEBDOMADAIRE
# ==========================================
slide12 = prs.slides.add_slide(slide_layout)
apply_background(slide12, COLOR_DARK)

# Title
tx_t = slide12.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0))
tf_t = tx_t.text_frame; tf_t.word_wrap = True
p_cat12 = tf_t.paragraphs[0]
p_cat12.text = "10 · PLAN D'ATTAQUE"
p_cat12.font.name = "Arial"; p_cat12.font.size = Pt(11); p_cat12.font.bold = True
p_cat12.font.color.rgb = COLOR_GOLD; p_cat12.space_after = Pt(4)

p_main12 = tf_t.add_paragraph()
p_main12.text = "La Routine du Succès — 5 Étapes Actionnables"
p_main12.font.name = "Arial"; p_main12.font.size = Pt(24); p_main12.font.bold = True
p_main12.font.color.rgb = COLOR_WHITE

steps12 = [
    (COLOR_BLUE, "01", "CIBLEZ",
     "Identifiez 50 prospects CVC RGE par semaine sans partenaire exclusif sur leur secteur.",
     "LinkedIn Sales Nav · Societe.com · Google Maps (recherche artisans CVC + ville)"),
    (COLOR_EMERALD, "02", "DIAGNOSTIQUEZ",
     "Pour chaque prospect, effectuez le tri-check rapide avant tout contact.",
     "site:son-site.com · Note Google Maps · Marques posées (Daikin, Mitsu, Fujitsu)"),
    (COLOR_ORANGE, "03", "CONTACTEZ",
     "Lancez la séquence mixte sur 5 jours (Social → Phoning → DM → Urgence).",
     "Instagram story-réaction · LinkedIn trigger · Phoning sujet 'gainable' · DM relance"),
    (COLOR_GOLD, "04", "DÉMONTREZ",
     "Faites la démo de 5 min en live avec l'effet WoW de l'écriture IA.",
     "IA article live · Carte exclusivité · Recherche Google en direct sous ses yeux"),
    (COLOR_RED, "05", "CLOSEZ",
     "Envoyez le lien Stripe avec le code promo EXCLU10 (-10%). Restez en ligne pendant la saisie.",
     "765 € HT au lieu de 850 € HT · Code promo = accélérateur de décision immédiate"),
]

for i, (color, num, label, desc, tools) in enumerate(steps12):
    x = 0.45 + i * 2.5
    y = 1.75

    # Main colored card
    card = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(2.3), Inches(5.5))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(40, 58, 79)
    card.line.color.rgb = color; card.line.width = Pt(2)

    # Number badge
    badge = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + 0.75), Inches(y + 0.2), Inches(0.8), Inches(0.55))
    badge.fill.solid(); badge.fill.fore_color.rgb = color; badge.line.color.rgb = color

    tx_n2 = slide12.shapes.add_textbox(Inches(x + 0.8), Inches(y + 0.25), Inches(0.7), Inches(0.4))
    tf_n2 = tx_n2.text_frame
    p_n2 = tf_n2.paragraphs[0]
    p_n2.text = num; p_n2.font.name = "Arial"; p_n2.font.size = Pt(14); p_n2.font.bold = True
    p_n2.font.color.rgb = COLOR_WHITE; p_n2.alignment = PP_ALIGN.CENTER

    # Label
    tx_label = slide12.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.85), Inches(2.1), Inches(0.55))
    tf_lbl = tx_label.text_frame
    p_lbl = tf_lbl.paragraphs[0]
    p_lbl.text = label
    p_lbl.font.name = "Arial"; p_lbl.font.size = Pt(14); p_lbl.font.bold = True
    p_lbl.font.color.rgb = color; p_lbl.alignment = PP_ALIGN.CENTER

    # Description
    tx_desc12 = slide12.shapes.add_textbox(Inches(x + 0.15), Inches(y + 1.5), Inches(2.0), Inches(2.5))
    tf_d12 = tx_desc12.text_frame; tf_d12.word_wrap = True
    p_d12 = tf_d12.paragraphs[0]
    p_d12.text = desc
    p_d12.font.name = "Arial"; p_d12.font.size = Pt(9.5)
    p_d12.font.color.rgb = COLOR_WHITE; p_d12.space_after = Pt(10)

    p_tools12 = tf_d12.add_paragraph()
    p_tools12.text = tools
    p_tools12.font.name = "Arial"; p_tools12.font.size = Pt(8.5)
    p_tools12.font.color.rgb = color; p_tools12.font.bold = True

# Bottom CTA strip
cta_strip = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(7.0), Inches(13.333), Inches(0.5))
cta_strip.fill.solid(); cta_strip.fill.fore_color.rgb = COLOR_GOLD
cta_strip.line.color.rgb = COLOR_GOLD

tx_cta = slide12.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.333), Inches(0.4))
tf_cta = tx_cta.text_frame
p_cta = tf_cta.paragraphs[0]
p_cta.text = "Gainable.fr — contact@gainable.fr — Code promo : EXCLU10 — 0 revente de contact — Le client vous choisit — 0% commission"
p_cta.font.name = "Arial"; p_cta.font.size = Pt(10); p_cta.font.bold = True
p_cta.font.color.rgb = COLOR_DARK; p_cta.alignment = PP_ALIGN.CENTER


# ==========================================
# SLIDE 13: COMPARATIF PLATEFORMES
# ==========================================
slide13 = prs.slides.add_slide(slide_layout)
apply_background(slide13, COLOR_LIGHT)
add_header(slide13, "Comparatif — Gainable.fr vs Bilik, Travaux.com, Selocal, BNI, PagesJaunes", "11 · POSITIONNEMENT CONCURRENTIEL")

# Intro badge
add_section_badge(slide13, 0.5, 1.62, 11.5, 0.32,
    "Maîtrisez chaque concurrent pour positionner Gainable.fr sans hésitation", COLOR_DARK)

# Table headers
headers = ["Critère", "Gainable.fr ✦", "Bilik", "Travaux.com", "Selocal", "BNI", "PagesJaunes"]
header_colors = [COLOR_DARK, COLOR_BLUE, COLOR_TEXT_SEC, COLOR_TEXT_SEC, COLOR_TEXT_SEC, COLOR_TEXT_SEC, COLOR_TEXT_SEC]
col_widths = [2.5, 2.0, 1.5, 1.6, 1.5, 1.3, 1.6]
col_starts = [0.3]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w + 0.05)

# Header row
header_bg = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(2.05), Inches(12.5), Inches(0.38))
header_bg.fill.solid(); header_bg.fill.fore_color.rgb = COLOR_DARK
header_bg.line.color.rgb = COLOR_DARK

for ci, (hdr, cx) in enumerate(zip(headers, col_starts)):
    tx_h = slide13.shapes.add_textbox(Inches(cx + 0.06), Inches(2.1), Inches(col_widths[ci] - 0.1), Inches(0.35))
    p_h = tx_h.text_frame.paragraphs[0]
    p_h.text = hdr
    p_h.font.name = "Arial"; p_h.font.size = Pt(8.5); p_h.font.bold = True
    p_h.font.color.rgb = COLOR_GOLD if ci == 1 else COLOR_WHITE
    p_h.alignment = PP_ALIGN.CENTER

# Data rows
rows_data = [
    ("Spécialisation CVC",      "✅ 100% CVC",             "⚠️ BTP gén.",     "⚠️ Tous travaux",   "⚠️ Tous métiers",  "❌ Hors sujet",   "❌ Généraliste"),
    ("Modèle de lead",          "✅ Visibilité directe",    "✅ Mise en rel.",  "❌ Lead x3-5 conc.", "⚠️ Abon.+leads",   "✅ Recommanda.",   "❌ Lead revendu"),
    ("Prix annuel",             "✅ 850 € HT fixe",         "⚠️ % commission", "⚠️ 40-80€/lead",    "⚠️ Variable",      "❌ 1500-2000€+",  "⚠️ Variable"),
    ("Commission chantier",     "✅ 0% commission",         "❌ 5-15% / vente","✅ 0%",              "✅ 0%",             "✅ 0%",            "✅ 0%"),
    ("SEO & Visibilité Google", "✅ +58k pages indexées",   "⚠️ Quelques pgs", "✅ Forte autorité",  "⚠️ Faible",        "❌ 0 SEO",        "⚠️ Vieillissant"),
    ("0 revente de contact",     "✅ Client vous choisit",   "✅ Contact direct",  "❌ Lead x3-5 conc.", "⚠️ Variable",    "✅ Recommandation","❌ Lead revendu"),
    ("IA SEO intégrée",         "✅ IA native 2 min",       "❌ Non",          "❌ Non",             "❌ Non",            "❌ Non",           "❌ Non"),
    ("Complémentaire ?",        "—",                       "✅ Oui",          "⚠️ Éviter",         "⚠️ Possible",      "✅ Très bien !",   "⚠️ Faible valeur"),
]

row_colors_alt = [COLOR_LIGHT, COLOR_WHITE]
for ri, row in enumerate(rows_data):
    y_row = 2.48 + ri * 0.54
    row_bg = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(y_row), Inches(12.5), Inches(0.52))
    row_bg.fill.solid(); row_bg.fill.fore_color.rgb = row_colors_alt[ri % 2]
    row_bg.line.color.rgb = COLOR_LINE

    for ci, (cell, cx) in enumerate(zip(row, col_starts)):
        tx_c = slide13.shapes.add_textbox(Inches(cx + 0.05), Inches(y_row + 0.08), Inches(col_widths[ci] - 0.1), Inches(0.4))
        p_c = tx_c.text_frame.paragraphs[0]
        p_c.text = cell
        p_c.font.name = "Arial"
        p_c.font.size = Pt(8.5) if ci == 0 else Pt(8)
        p_c.font.bold = (ci == 0 or ci == 1)
        if ci == 1:
            p_c.font.color.rgb = COLOR_BLUE
        elif cell.startswith("✅"):
            p_c.font.color.rgb = COLOR_EMERALD
        elif cell.startswith("❌"):
            p_c.font.color.rgb = COLOR_RED
        elif cell.startswith("⚠"):
            p_c.font.color.rgb = COLOR_ORANGE
        else:
            p_c.font.color.rgb = COLOR_TEXT_PRI
        p_c.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT

# Bottom tip box
tip_box = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(6.97), Inches(12.5), Inches(0.42))
tip_box.fill.solid(); tip_box.fill.fore_color.rgb = RGBColor(255, 251, 235)
tip_box.line.color.rgb = COLOR_GOLD; tip_box.line.width = Pt(1.5)
tx_tip13 = slide13.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.1), Inches(0.38))
p_tip13 = tx_tip13.text_frame.paragraphs[0]
p_tip13.text = ("💡 BNI = complément parfait (réseau humain chaud). Travaux.com = anti-modèle à citer : "
                "« Vous connaissez le stress d'appeler en 3 min avant 4 concurrents pour brader vos marges ? Chez nous, le client vous appelle VOUS seul. »")
p_tip13.font.name = "Arial"; p_tip13.font.size = Pt(9); p_tip13.font.bold = True
p_tip13.font.color.rgb = COLOR_DARK

# ==========================================
# SLIDE 14: ÉCOSYSTÈME DE VISIBILITÉ
# ==========================================
slide14 = prs.slides.add_slide(slide_layout)
apply_background(slide14, COLOR_LIGHT)
add_header(slide14, "L'Écosystème de Visibilité Gainable.fr — Vue Globale", "12 · MAILLAGE DIGITAL")

# Center hub
hub = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.65), Inches(3.7), Inches(1.1))
hub.fill.solid(); hub.fill.fore_color.rgb = COLOR_DARK
hub.line.color.rgb = COLOR_GOLD; hub.line.width = Pt(2.5)
tx_hub = slide14.shapes.add_textbox(Inches(4.9), Inches(1.75), Inches(3.5), Inches(0.9))
tf_hub = tx_hub.text_frame
p_hub = tf_hub.paragraphs[0]
p_hub.text = "Gainable.fr"
p_hub.font.name = "Arial"; p_hub.font.size = Pt(18); p_hub.font.bold = True
p_hub.font.color.rgb = COLOR_GOLD; p_hub.alignment = PP_ALIGN.CENTER
p_hub2 = tf_hub.add_paragraph()
p_hub2.text = "Hub Central • +58k pages • 3 pays • 100% CVC"
p_hub2.font.name = "Arial"; p_hub2.font.size = Pt(8.5)
p_hub2.font.color.rgb = COLOR_WHITE; p_hub2.alignment = PP_ALIGN.CENTER

# 6 levers
levers = [
    (0.35, 2.9,  COLOR_BLUE,    "🔍 SEO Google",          ["Articles géolocalisés", "Mots-clés CVC longue traîne", "Rich Snippets JSON-LD", "Autorité domaine collective"]),
    (2.65, 2.9,  COLOR_EMERALD, "🌐 Site Officiel",        ["Fiche Pro optimisée", "Photos de chantiers", "Certifications RGE", "Coordonnées directes (0%)"] ),
    (5.1,  4.95, COLOR_ORANGE,  "📱 Réseaux Sociaux",      ["Instagram → leads DM", "LinkedIn → B2B", "Signaux sociaux Google", "Bouche-à-oreille digital"]),
    (7.65, 2.9,  COLOR_RED,     "📍 Google Maps / GMB",   ["Backlink Gainable→GMB", "Avis croisés", "Pack Local Google", "Requêtes 'près de moi'"] ),
    (10.0, 2.9,  COLOR_GOLD,    "✍️ Contenu IA",           ["10 articles/mois", "120/an par expert", "FAQ JSON-LD", "Cocon sémantique auto"]),
    (5.1,  1.65, RGBColor(139, 92, 246), "🤖 IA & AEO",  ["ChatGPT, Gemini, Perplexity", "Réponses directes IA (AEO)", "Schema.org complet", "Position 0 générative"]),
]

for lx, ly, lc, ltitle, litems in levers:
    # Skip hub-overlapping levers (IA is above hub)
    card14 = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(lx), Inches(ly), Inches(2.15), Inches(2.2))
    card14.fill.solid(); card14.fill.fore_color.rgb = COLOR_WHITE
    card14.line.color.rgb = lc; card14.line.width = Pt(1.5)

    top14 = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(lx), Inches(ly), Inches(2.15), Inches(0.07))
    top14.fill.solid(); top14.fill.fore_color.rgb = lc; top14.line.color.rgb = lc

    tx14 = slide14.shapes.add_textbox(Inches(lx + 0.1), Inches(ly + 0.1), Inches(1.95), Inches(2.05))
    tf14 = tx14.text_frame; tf14.word_wrap = True

    p14_t = tf14.paragraphs[0]
    p14_t.text = ltitle
    p14_t.font.name = "Arial"; p14_t.font.size = Pt(10); p14_t.font.bold = True
    p14_t.font.color.rgb = lc; p14_t.space_after = Pt(4)

    for item14 in litems:
        p14_i = tf14.add_paragraph()
        p14_i.text = f"• {item14}"
        p14_i.font.name = "Arial"; p14_i.font.size = Pt(8.5)
        p14_i.font.color.rgb = COLOR_TEXT_SEC; p14_i.space_after = Pt(2)

# Effect banner
effect_bar = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.35), Inches(7.0), Inches(12.5), Inches(0.38))
effect_bar.fill.solid(); effect_bar.fill.fore_color.rgb = RGBColor(239, 246, 255)
effect_bar.line.color.rgb = COLOR_BLUE

tx_effect = slide14.shapes.add_textbox(Inches(0.5), Inches(7.03), Inches(12.2), Inches(0.35))
p_eff = tx_effect.text_frame.paragraphs[0]
p_eff.text = ("⚡ Effet de Maillage = Amplification Exponentielle : Article publié → Indexé Google → "
              "Partagé Instagram → Cité par l'IA → Visible sur Maps → L'artisan reçoit des appels sans rien faire.")
p_eff.font.name = "Arial"; p_eff.font.size = Pt(9); p_eff.font.bold = True
p_eff.font.color.rgb = COLOR_BLUE

# ==========================================
# SLIDE 15: LES 3 PHASES DU MAILLAGE
# ==========================================
slide15 = prs.slides.add_slide(slide_layout)
apply_background(slide15, COLOR_LIGHT)
add_header(slide15, "Les 3 Phases du Maillage — De l'Ancrage à la Dominance Locale", "13 · STRATÉGIE MAILLAGE")

phases = [
    (COLOR_BLUE, "📍", "Phase 1 — Ancrage", "Mois 1 à 3",
     ["Fiche Pro complète : description, certifications, marques posées",
      "3 photos minimum de chantiers récents (intérieur + technique)",
      "Coordonnées directes vérifiées (téléphone + email)",
      "Zone géographique définie avec rayon d'intervention",
      "Premier article SEO optimisé sur la ville principale",
      "Lien depuis GMB → Gainable.fr activé"],
     "Google commence à crawler et indexer la fiche — les fondations sont posées."),
    (COLOR_EMERALD, "🕸️", "Phase 2 — Maillage", "Mois 3 à 6",
     ["8 à 10 articles/mois sur villes + technologies (gainable + multi-split…)",
      "FAQ structurées JSON-LD intégrées pour Featured Snippets",
      "Photos optimisées avec balises Alt localisées",
      "Partages Instagram & LinkedIn des articles publiés",
      "Avis Google croisés avec la fiche Gainable.fr",
      "Backlinks depuis partenaires locaux (électriciens, plombiers…)"],
     "Premières positions sur requêtes longue traîne locales — premiers contacts entrants."),
    (COLOR_GOLD, "🏆", "Phase 3 — Dominance", "Mois 6 et au-delà",
     ["Cocon sémantique complet (50+ articles sur la zone)",
      "Featured Snippets capturés via FAQ JSON-LD",
      "Réponses IA (ChatGPT, Gemini) qui citent l'artisan",
      "Backlinks naturels depuis blogs & médias locaux",
      "Page 1 sur requêtes principales de la zone",
      "2 à 5 contacts qualifiés par mois en entrant"],
     "Leads entrants réguliers — l'artisan reçoit des appels sans prospecter."),
]

for pi, (color, icon, phase_title, timing, actions, result) in enumerate(phases):
    x = 0.4 + pi * 4.3
    y = 1.68

    # Main card
    card15 = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.1), Inches(5.6))
    card15.fill.solid(); card15.fill.fore_color.rgb = COLOR_WHITE
    card15.line.color.rgb = color; card15.line.width = Pt(2.0)

    # Top banner
    banner15 = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(4.1), Inches(1.0))
    banner15.fill.solid(); banner15.fill.fore_color.rgb = color
    banner15.line.color.rgb = color

    # Banner text
    tx_b15 = slide15.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.07), Inches(3.8), Inches(0.9))
    tf_b15 = tx_b15.text_frame
    p_b15_icon = tf_b15.paragraphs[0]
    p_b15_icon.text = f"{icon}  {phase_title}"
    p_b15_icon.font.name = "Arial"; p_b15_icon.font.size = Pt(13); p_b15_icon.font.bold = True
    p_b15_icon.font.color.rgb = COLOR_WHITE; p_b15_icon.space_after = Pt(2)
    p_b15_timing = tf_b15.add_paragraph()
    p_b15_timing.text = timing
    p_b15_timing.font.name = "Arial"; p_b15_timing.font.size = Pt(10)
    p_b15_timing.font.color.rgb = RGBColor(255, 255, 200)

    # Actions
    tx_a15 = slide15.shapes.add_textbox(Inches(x + 0.15), Inches(y + 1.08), Inches(3.8), Inches(3.5))
    tf_a15 = tx_a15.text_frame; tf_a15.word_wrap = True
    first15 = True
    for act in actions:
        if first15:
            p_a15 = tf_a15.paragraphs[0]
            first15 = False
        else:
            p_a15 = tf_a15.add_paragraph()
        p_a15.text = f"✓  {act}"
        p_a15.font.name = "Arial"; p_a15.font.size = Pt(9.5)
        p_a15.font.color.rgb = COLOR_TEXT_PRI; p_a15.space_after = Pt(4)

    # Result box
    res_box = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + 0.1), Inches(y + 4.55), Inches(3.9), Inches(0.9))
    res_box.fill.solid(); res_box.fill.fore_color.rgb = RGBColor(241, 245, 249)
    res_box.line.color.rgb = color; res_box.line.width = Pt(1)
    tx_r15 = slide15.shapes.add_textbox(Inches(x + 0.2), Inches(y + 4.62), Inches(3.7), Inches(0.78))
    tf_r15 = tx_r15.text_frame; tf_r15.word_wrap = True
    p_r15 = tf_r15.paragraphs[0]
    p_r15.text = f"🎯 {result}"
    p_r15.font.name = "Arial"; p_r15.font.size = Pt(9); p_r15.font.bold = True
    p_r15.font.color.rgb = color

# ==========================================
# SLIDE 16: MINDSET & ÉTAT D'ESPRIT
# ==========================================
slide16 = prs.slides.add_slide(slide_layout)
apply_background(slide16, COLOR_DARK)

# Title
tx_t16 = slide16.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.7), Inches(1.1))
tf_t16 = tx_t16.text_frame; tf_t16.word_wrap = True
p_cat16 = tf_t16.paragraphs[0]
p_cat16.text = "14 · MINDSET & ÉTAT D'ESPRIT"
p_cat16.font.name = "Arial"; p_cat16.font.size = Pt(10); p_cat16.font.bold = True
p_cat16.font.color.rgb = COLOR_GOLD; p_cat16.space_after = Pt(4)
p_title16 = tf_t16.add_paragraph()
p_title16.text = "La Conviction qui Fait la Différence — Leche Forge et Rend Plus Fort"
p_title16.font.name = "Arial"; p_title16.font.size = Pt(20); p_title16.font.bold = True
p_title16.font.color.rgb = COLOR_WHITE

mindset_cards = [
    (COLOR_RED,     "🔥", "Le Refus n'est Pas un Échec",
     "Un 'non' = 'pas encore'. Sur 50 prospects : 10 décrochent, 3 sont intéressés, 1-2 signent. Ce ratio est NORMAL. Chaque refus rapproche la signature.",
     "Fixez-vous un quota de refus : 15 'non'/jour = vous avez bien travaillé."),
    (COLOR_BLUE,    "🛡️", "Croire en son Produit — Conviction Absolue",
     "Gainable.fr crée de la VRAIE valeur. L'IA rédige. Les pages sont indexées. Quand vous présentez, vous n'êtes pas en train de demander de l'argent — vous offrez une opportunité.",
     "Relisez 1 témoignage de partenaire satisfait chaque matin avant vos appels."),
    (COLOR_EMERALD, "📈", "La Règle des 50 Touches",
     "Un prospect signe après 5 à 7 touches en moyenne. La plupart des commerciaux abandonnent après 1-2. En suivant la cadence de 5 jours, vous allez là où les autres ne vont pas.",
     "Tracez CHAQUE action dans le CRM. Ne classez 'refus définitif' qu'après J5."),
    (COLOR_GOLD,    "⚔️", "Le Rituel du Guerrier",
     "Votre état d'esprit se détecte en 3 secondes à la voix. Un commercial fatigué transmet sa fatigue. La préparation mentale est un avantage compétitif concret.",
     "2 min de respiration avant les appels. Debout pendant le phoning. Sourire — ça s'entend."),
    (COLOR_ORANGE,  "🌱", "Planter et Ne Jamais Abandonner",
     "Un artisan qui dit 'non' en janvier peut signer en mars quand il voit son concurrent sur Google. Votre travail aujourd'hui génère des résultats dans 90 jours.",
     "Remontez dans vos 'refus de 3 mois' — c'est votre meilleure base de rappel."),
    (COLOR_RED,     "🏆", "Célébrer Chaque Victoire",
     "Une vente à 850 € HT, c'est entre 85 € et 127 € de commission. Mais c'est aussi un artisan dont la vie va changer. Ce que l'on célèbre, on le reproduit.",
     "Notez chaque vente dans un carnet visible. La liste grandit, ça motive."),
]

col_pos16 = [(0.35, 6.3), (6.85, 6.3)]
row_y16 = [1.6, 3.3, 5.0]

for mi, (color, icon, title, content, tip) in enumerate(mindset_cards):
    col = mi % 2
    row = mi // 2
    x = col_pos16[col][0]
    y = row_y16[row]
    w = col_pos16[col][1]

    card16 = slide16.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.6))
    card16.fill.solid(); card16.fill.fore_color.rgb = RGBColor(40, 58, 79)
    card16.line.color.rgb = color; card16.line.width = Pt(1.5)

    left16 = slide16.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(1.6))
    left16.fill.solid(); left16.fill.fore_color.rgb = color; left16.line.color.rgb = color

    tx16 = slide16.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.1), Inches(w - 0.28), Inches(1.45))
    tf16 = tx16.text_frame; tf16.word_wrap = True

    p16_t = tf16.paragraphs[0]
    p16_t.text = f"{icon}  {title}"
    p16_t.font.name = "Arial"; p16_t.font.size = Pt(11); p16_t.font.bold = True
    p16_t.font.color.rgb = color; p16_t.space_after = Pt(3)

    p16_c = tf16.add_paragraph()
    p16_c.text = content
    p16_c.font.name = "Arial"; p16_c.font.size = Pt(8.5)
    p16_c.font.color.rgb = COLOR_WHITE; p16_c.space_after = Pt(3)

    p16_tip = tf16.add_paragraph()
    p16_tip.text = f"→ {tip}"
    p16_tip.font.name = "Arial"; p16_tip.font.size = Pt(8.5); p16_tip.font.bold = True
    p16_tip.font.color.rgb = color

# ==========================================
# SLIDE 17: MARATHON — LIGNE DU TEMPS
# ==========================================
slide17 = prs.slides.add_slide(slide_layout)
apply_background(slide17, COLOR_DARK)

# Title
tx_t17 = slide17.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.7), Inches(1.1))
tf_t17 = tx_t17.text_frame; tf_t17.word_wrap = True
p_cat17 = tf_t17.paragraphs[0]
p_cat17.text = "15 · VISION LONG TERME"
p_cat17.font.name = "Arial"; p_cat17.font.size = Pt(10); p_cat17.font.bold = True
p_cat17.font.color.rgb = COLOR_GOLD; p_cat17.space_after = Pt(4)
p_title17 = tf_t17.add_paragraph()
p_title17.text = "C'est un Marathon, Pas un Sprint — La Ligne du Temps SEO"
p_title17.font.name = "Arial"; p_title17.font.size = Pt(20); p_title17.font.bold = True
p_title17.font.color.rgb = COLOR_WHITE

# Subtitle
tx_sub17 = slide17.shapes.add_textbox(Inches(0.8), Inches(1.52), Inches(11.7), Inches(0.4))
p_sub17 = tx_sub17.text_frame.paragraphs[0]
p_sub17.text = "La visibilité Google ne se construit pas en une nuit. Plus l'artisan est actif sur Gainable.fr, plus son autorité grandit et plus il domine sa zone."
p_sub17.font.name = "Arial"; p_sub17.font.size = Pt(10)
p_sub17.font.color.rgb = COLOR_TEXT_SEC

# Timeline horizontal bar
bar = slide17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.15), Inches(12.333), Inches(0.1))
bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(55, 65, 81); bar.line.color.rgb = RGBColor(55, 65, 81)

# Progress fill (showing m1-3 as "now")
prog = slide17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.15), Inches(3.1), Inches(0.1))
prog.fill.solid(); prog.fill.fore_color.rgb = COLOR_BLUE; prog.line.color.rgb = COLOR_BLUE

# Timeline milestones
milestones = [
    (0.5,   COLOR_BLUE,    "🌱", "Mois 1-3",  "FONDATIONS",
     ["Fiche indexée par Google", "Premiers articles publiés", "Google commence à voir l'artisan"]),
    (3.6,   COLOR_EMERALD, "📈", "Mois 3-6",  "MONTÉE EN PUISSANCE",
     ["Page 2-3 sur requêtes locales", "Premiers contacts entrants", "Confiance Google augmente"]),
    (6.7,   COLOR_GOLD,    "📞", "Mois 6-12", "LEADS RÉGULIERS",
     ["Page 1 sur requêtes clés", "2-5 contacts qualifiés/mois", "ROI pleinement atteint"]),
    (9.8,   COLOR_RED,     "🏆", "12 mois+",  "DOMINANCE LOCALE",
     ["Référent SEO sur la zone", "Concurrent difficile à déloger", "Chantiers premium en priorité"]),
]

for mx, mc, micon, mperiod, mlabel, mpoints in milestones:
    # Dot on timeline
    dot = slide17.shapes.add_shape(MSO_SHAPE.OVAL, Inches(mx + 1.15), Inches(3.06), Inches(0.27), Inches(0.27))
    dot.fill.solid(); dot.fill.fore_color.rgb = mc; dot.line.color.rgb = mc

    # Card below
    card17 = slide17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(mx), Inches(3.5), Inches(3.0), Inches(3.5))
    card17.fill.solid(); card17.fill.fore_color.rgb = RGBColor(40, 58, 79)
    card17.line.color.rgb = mc; card17.line.width = Pt(2)

    top17 = slide17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(mx), Inches(3.5), Inches(3.0), Inches(0.07))
    top17.fill.solid(); top17.fill.fore_color.rgb = mc; top17.line.color.rgb = mc

    tx17 = slide17.shapes.add_textbox(Inches(mx + 0.15), Inches(3.57), Inches(2.7), Inches(3.38))
    tf17 = tx17.text_frame; tf17.word_wrap = True

    p17_icon = tf17.paragraphs[0]
    p17_icon.text = f"{micon}  {mperiod}"
    p17_icon.font.name = "Arial"; p17_icon.font.size = Pt(9); p17_icon.font.bold = True
    p17_icon.font.color.rgb = mc; p17_icon.space_after = Pt(2)

    p17_label = tf17.add_paragraph()
    p17_label.text = mlabel
    p17_label.font.name = "Arial"; p17_label.font.size = Pt(11); p17_label.font.bold = True
    p17_label.font.color.rgb = COLOR_WHITE; p17_label.space_after = Pt(8)

    for pt17 in mpoints:
        p17_pt = tf17.add_paragraph()
        p17_pt.text = f"• {pt17}"
        p17_pt.font.name = "Arial"; p17_pt.font.size = Pt(9.5)
        p17_pt.font.color.rgb = COLOR_TEXT_SEC; p17_pt.space_after = Pt(3)

# Key argument box
arg_box = slide17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.38))
arg_box.fill.solid(); arg_box.fill.fore_color.rgb = RGBColor(27, 37, 50)
arg_box.line.color.rgb = COLOR_GOLD; arg_box.line.width = Pt(1.5)
tx_arg17 = slide17.shapes.add_textbox(Inches(0.7), Inches(7.03), Inches(12.0), Inches(0.35))
p_arg17 = tx_arg17.text_frame.paragraphs[0]
p_arg17.text = ("💡 Argument clé : « Le meilleur moment pour planter un arbre était il y a 1 an. Le deuxième meilleur, c'est AUJOURD'HUI. "
                "Vos concurrents déjà inscrits récoltent. Chaque mois d'attente = 1 mois de retard. »")
p_arg17.font.name = "Arial"; p_arg17.font.size = Pt(9.5); p_arg17.font.bold = True
p_arg17.font.color.rgb = COLOR_GOLD

# ==========================================
# SAVE
# ==========================================
output_file = os.path.join(project_root, "playbook_cvc_gainable_v2.pptx")
try:
    prs.save(output_file)
    print("OK Playbook sauvegarde : " + output_file)
except PermissionError:
    output_file = os.path.join(project_root, "playbook_cvc_gainable_v2_new.pptx")
    prs.save(output_file)
    print("OK Playbook sauvegarde (fallback) : " + output_file)


