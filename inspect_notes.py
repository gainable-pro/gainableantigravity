from pptx import Presentation

prs = Presentation('gainable_presentation_final.pptx')

for i, slide in enumerate(prs.slides):
    notes_slide = slide.notes_slide
    if notes_slide and notes_slide.notes_text_frame and notes_slide.notes_text_frame.text.strip():
        print(f"Slide {i+1} Notes: {notes_slide.notes_text_frame.text}")
