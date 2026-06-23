import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Scheme Definitions (matching Gainable.fr brand)
COLOR_DARK = RGBColor(31, 45, 61)       # #1F2D3D (Dark slate / navy)
COLOR_LIGHT = RGBColor(248, 250, 252)   # #F8FAFC (Slate-50)
COLOR_GOLD = RGBColor(213, 155, 43)     # #D59B2B (Gainable Gold)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_TEXT_PRIMARY = RGBColor(31, 45, 61)
COLOR_TEXT_SECONDARY = RGBColor(100, 116, 139) # Slate-500
COLOR_LINE = RGBColor(226, 232, 240)    # Slate-200
COLOR_EMERALD = RGBColor(16, 185, 129)  # #10B981

project_root = r"C:\Users\gmaro\.gemini\antigravity-ide\scratch\gainable-fr"
path_logo = os.path.join(project_root, "public", "logo.png")
path_photo_meeting = os.path.join(project_root, "public", "sales_meeting.png")
path_photo_handshake = os.path.join(project_root, "public", "sales_handshake.png")

def apply_background(slide, color):
    """Fills the slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, category_text="FORMATION COMMERCIALE"):
    """Adds a standard premium header to light-background slides."""
    # Category prefix (small gold text)
    txBox_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
    tf_cat = txBox_cat.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_GOLD
    
    # Title (large navy text)
    txBox_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf_title = txBox_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Arial"
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_PRIMARY
    
    # Subtle underline border for the header area
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.733), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_LINE
    line.line.color.rgb = COLOR_LINE

def format_bold_lead(paragraph, bold_text, normal_text, size=Pt(12.5), font_color=COLOR_TEXT_PRIMARY):
    """Formats a paragraph with bold prefix and regular body text."""
    paragraph.font.name = "Arial"
    paragraph.font.size = size
    paragraph.space_after = Pt(8)
    
    run_bold = paragraph.add_run()
    run_bold.text = bold_text
    run_bold.font.bold = True
    run_bold.font.color.rgb = font_color
    
    run_normal = paragraph.add_run()
    run_normal.text = normal_text
    run_normal.font.bold = False
    run_normal.font.color.rgb = COLOR_TEXT_SECONDARY

# ==========================================
# SLIDE 1: Title Slide (Dark Theme)
# ==========================================
slide_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(slide_layout)
apply_background(slide1, COLOR_DARK)

# Main Title Panel
tx_title = slide1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(6.5), Inches(4.5))
tf_title = tx_title.text_frame
tf_title.word_wrap = True

# Logo in Title Panel
if os.path.exists(path_logo):
    slide1.shapes.add_picture(path_logo, Inches(0.8), Inches(1.0), width=Inches(2.6), height=Inches(0.78))

# Right Side: Sales Meeting Photo
if os.path.exists(path_photo_meeting):
    slide1.shapes.add_picture(path_photo_meeting, Inches(7.333), Inches(0), width=Inches(6.0), height=Inches(7.5))

p_t = tf_title.paragraphs[0]
p_t.text = "GUIDE COMMERCIAL & DE VENTE"
p_t.font.name = "Arial"
p_t.font.size = Pt(40)
p_t.font.bold = True
p_t.font.color.rgb = COLOR_WHITE
p_t.space_after = Pt(6)

p_t_sub = tf_title.add_paragraph()
p_t_sub.text = "Processus, Écoute Active & Traitement des Objections Artisans"
p_t_sub.font.name = "Arial"
p_t_sub.font.size = Pt(20)
p_t_sub.font.bold = True
p_t_sub.font.color.rgb = COLOR_GOLD
p_t_sub.space_after = Pt(20)

p_t_desc = tf_title.add_paragraph()
p_t_desc.text = "Destiné aux consultants commerciaux Gainable.fr.\nApprenez à cibler le bon client, écouter ses besoins réels et surmonter chaque objection."
p_t_desc.font.name = "Arial"
p_t_desc.font.size = Pt(13)
p_t_desc.font.color.rgb = COLOR_WHITE
p_t_desc.space_after = Pt(30)

p_t_foot = tf_title.add_paragraph()
p_t_foot.text = "Document Interne Confidentiel - Gainable.fr 2026"
p_t_foot.font.name = "Arial"
p_t_foot.font.size = Pt(11)
p_t_foot.font.color.rgb = COLOR_TEXT_SECONDARY

# ==========================================
# SLIDE 2: Problématiques des Artisans (HVAC / Diag / BE)
# ==========================================
slide2 = prs.slides.add_slide(slide_layout)
apply_background(slide2, COLOR_LIGHT)
add_header(slide2, "Comprendre son prospect : Les problématiques des artisans", "CONNAISSANCE MARCHÉ")

tx_arg2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf_arg2 = tx_arg2.text_frame
tf_arg2.word_wrap = True

p_intro2 = tf_arg2.paragraphs[0]
p_intro2.text = "Pour vendre efficacement, vous devez résonner avec les douleurs réelles rencontrées par les experts CVC, diagnostiqueurs et bureaux d'études au quotidien :"
p_intro2.font.bold = True
p_intro2.font.size = Pt(14)
p_intro2.font.color.rgb = COLOR_DARK
p_intro2.space_after = Pt(15)

p2_b1 = tf_arg2.add_paragraph()
format_bold_lead(p2_b1, "1. La revente multiple de leads (Le grand fléau) : ", "Les plateformes de devis traditionnelles revendent le même prospect à 3, 4 ou 5 concurrents. Cela déclenche une guerre des prix agressive vers le bas, détruisant leurs marges et leur temps.")

p2_b2 = tf_arg2.add_paragraph()
format_bold_lead(p2_b2, "2. Le manque de temps et d'expertise digitale : ", "Un artisan passe ses journées sur les chantiers. Il n'a ni le temps ni les compétences pour optimiser son référencement Google, rédiger du contenu SEO ou configurer des campagnes publicitaires complexes.")

p2_b3 = tf_arg2.add_paragraph()
format_bold_lead(p2_b3, "3. Le coût publicitaire exorbitant (Google/Meta Ads) : ", "Gérer seul du Google Ads est risqué et coûteux. Sans expertise, ils dépensent des centaines d'euros pour de simples clics de curieux sans jamais signer de contrat.")

p2_b4 = tf_arg2.add_paragraph()
format_bold_lead(p2_b4, "4. L'isolement commercial : ", "Les TPE/PME du bâtiment n'ont pas les moyens des grands groupes pour s'imposer sur le web. Ils ont besoin d'un partenaire d'autorité globale pour les propulser localement.")

# ==========================================
# SLIDE 3: Détection des Besoins avant l'appel (La Préparation)
# ==========================================
slide3 = prs.slides.add_slide(slide_layout)
apply_background(slide3, COLOR_LIGHT)
add_header(slide3, "Processus : La détection des besoins avant d'appeler le client", "MÉTHODOLOGIE COMMERCIALE")

tx_s3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf_s3 = tx_s3.text_frame
tf_s3.word_wrap = True

ps3_intro = tf_s3.paragraphs[0]
ps3_intro.text = "Ne passez jamais d'appel à froid sans préparation. L'analyse préalable permet de personnaliser l'approche commerciale et d'établir une autorité immédiate :"
ps3_intro.font.bold = True
ps3_intro.font.size = Pt(14)
ps3_intro.font.color.rgb = COLOR_DARK
ps3_intro.space_after = Pt(15)

ps3_b1 = tf_s3.add_paragraph()
format_bold_lead(ps3_b1, "• Analyse de la présence web : ", "Recherchez le nom de l'entreprise sur Google. Ont-ils un site web ? Est-il responsive ? Ont-ils des avis Google My Business récents ? Si leur site est daté, c'est un point d'appui majeur pour notre solution de fiches pros optimisées.")

ps3_b2 = tf_s3.add_paragraph()
format_bold_lead(ps3_b2, "• Audit SEO de sa localité : ", "Faites une recherche rapide Google (ex: 'climatisation gainable + leur ville'). Regardez si l'artisan apparaît dans les premiers résultats. Montrez-lui que Gainable.fr y figure déjà, et qu'adhérer au réseau lui permet d'occuper cette place immédiatement.")

ps3_b3 = tf_s3.add_paragraph()
format_bold_lead(ps3_b3, "• Détection des budgets publicitaires : ", "Regardez s'ils achètent des mots-clés Google Ads (présence du badge 'Sponsorisé'). S'ils paient de la publicité payante Google, ils sont conscients de l'importance du digital mais paient probablement trop cher par lead par rapport à notre abonnement fixe.")

ps3_b4 = tf_s3.add_paragraph()
format_bold_lead(ps3_b4, "• Ciblage précis de sa typologie : ", "Ciblez le bon interlocuteur en vérifiant le code APE et le dirigeant. Proposez l'offre Expert CVC (850 € HT) aux installateurs, et Diag Immo (750 € HT) aux diagnostiqueurs. Adaptez le vocabulaire technique.")

# ==========================================
# SLIDE 4: Écoute Active & Empathie lors de l'appel
# ==========================================
slide4 = prs.slides.add_slide(slide_layout)
apply_background(slide4, COLOR_LIGHT)
add_header(slide4, "L'Appel : Pratiquer l'écoute active et l'empathie", "POSTURE COMMERCIALE")

tx_s4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf_s4 = tx_s4.text_frame
tf_s4.word_wrap = True

ps4_intro = tf_s4.paragraphs[0]
ps4_intro.text = "L'écoute active représente 70% de la réussite d'une vente. Ne parlez pas immédiatement de technique ou de tarif, laissez l'artisan exprimer ses frustrations :"
ps4_intro.font.bold = True
ps4_intro.font.size = Pt(14)
ps4_intro.font.color.rgb = COLOR_DARK
ps4_intro.space_after = Pt(15)

ps4_b1 = tf_s4.add_paragraph()
format_bold_lead(ps4_b1, "• Valider ses frustrations avec empathie : ", "S'il se plaint des plateformes de devis classiques ou de clients 'curieux', rebondissez : 'Je comprends tout à fait, c'est très frustrant de payer pour des leads partagés avec 4 concurrents et de ne pas être rappelé.' Cela crée une alliance.")

ps4_b2 = tf_s4.add_paragraph()
format_bold_lead(ps4_b2, "• Poser des questions ouvertes ciblées : ", "Posez des questions du type : 'Comment faites-vous aujourd'hui pour combler les périodes creuses ?' ou 'Combien vous coûte l'acquisition d'un nouveau client ?'. Laissez-le quantifier lui-même ses dépenses.")

ps4_b3 = tf_s4.add_paragraph()
format_bold_lead(ps4_b3, "• Reformuler pour montrer sa compréhension : ", "'Si je comprends bien, vous cherchez un canal qui vous apporte des clients en direct dans votre zone, sans intermédiaire, et sans que vos marges soient détruites ?' Le prospect doit valider par un 'Oui' franc.")

ps4_b4 = tf_s4.add_paragraph()
format_bold_lead(ps4_b4, "• Présenter Gainable.fr comme l'anti-plateforme : ", "Insistez sur la mise en relation directe (le client l'appelle directement, pas de revente), et l'abonnement annuel fixe (pas de commission au chantier).")

# ==========================================
# SLIDE 5: Traitement des Objections - Partie 1
# ==========================================
slide5 = prs.slides.add_slide(slide_layout)
apply_background(slide5, COLOR_LIGHT)
add_header(slide5, "Traitement des objections courantes (Partie 1)", "OBJECTIONS & CONTRE-ARGUMENTS")

tx_s5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf_s5 = tx_s5.text_frame
tf_s5.word_wrap = True

ps5_intro = tf_s5.paragraphs[0]
ps5_intro.text = "Soyez prêt à surmonter les deux barrières les plus fréquentes en réorientant sur le retour sur investissement (ROI) :"
ps5_intro.font.bold = True
ps5_intro.font.size = Pt(14)
ps5_intro.font.color.rgb = COLOR_DARK
ps5_intro.space_after = Pt(15)

ps5_b1 = tf_s5.add_paragraph()
format_bold_lead(ps5_b1, "Objection A : « C'est trop cher ! (850 € ou 750 € HT par an) »\n", 
                 "• Contre-argument 1 : Ramenez le tarif à l'échelle quotidienne : cela représente moins de 2,30 € HT par jour. C'est inférieur au prix d'un café !\n"
                 "• Contre-argument 2 : Comparez avec d'autres canaux : sur les plateformes de leads classiques, un seul lead climatisation qualifié coûte entre 40 et 80 € HT, qu'il soit signé ou non. Chez nous, l'adhésion est rentabilisée dès le premier chantier signé.\n"
                 "• Contre-argument 3 : Insistez sur le 0% commission : si le professionnel signe un chantier à 10 000 €, il ne nous doit rien de plus.")

ps5_b2 = tf_s5.add_paragraph()
format_bold_lead(ps5_b2, "Objection B : « J'ai déjà testé des plateformes de devis, ça ne marche pas, les leads sont mauvais. »\n", 
                 "• Contre-argument 1 : Expliquez notre différence fondamentale : nous ne sommes pas une plateforme de vente de leads. Nous sommes un annuaire de visibilité directe.\n"
                 "• Contre-argument 2 : Le parcours utilisateur : les particuliers nous trouvent via des recherches Google ultra-ciblées. Ils consultent la fiche de l'artisan, voient ses réalisations, et décident de l'appeler directement. Il n'y a aucun intermédiaire et aucune revente multiple.\n"
                 "• Contre-argument 3 : Des fiches d'intervention détaillées pour limiter les curieux.")

# ==========================================
# SLIDE 6: Traitement des Objections - Partie 2
# ==========================================
slide6 = prs.slides.add_slide(slide_layout)
apply_background(slide6, COLOR_LIGHT)
add_header(slide6, "Traitement des objections courantes (Partie 2)", "OBJECTIONS & CONTRE-ARGUMENTS")

tx_s6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf_s6 = tx_s6.text_frame
tf_s6.word_wrap = True

ps6_intro = tf_s6.paragraphs[0]
ps6_intro.text = "Savoir traiter le manque de temps ou l'autosatisfaction sur leur carnet de commandes actuel :"
ps6_intro.font.bold = True
ps6_intro.font.size = Pt(14)
ps6_intro.font.color.rgb = COLOR_DARK
ps6_intro.space_after = Pt(15)

ps6_b1 = tf_s6.add_paragraph()
format_bold_lead(ps6_b1, "Objection C : « Je n'ai pas le temps de rédiger des articles ou de m'occuper du SEO. »\n", 
                 "• Contre-argument 1 : L'Assistant Rédacteur IA intégré : l'artisan n'a rien à écrire lui-même. Notre Espace Pro intègre un rédacteur IA expert en génie climatique.\n"
                 "• Contre-argument 2 : Rédiger en 2 minutes : il saisit le nom de sa ville, décrit son dernier chantier en 3 mots ('installation bi-split Daikin'), et l'IA génère un article optimisé de 800 mots prêt à publier en un clic.\n"
                 "• Contre-argument 3 : C'est notre plateforme qui pousse la technique SEO, ils n'ont aucune configuration à faire.")

ps6_b2 = tf_s6.add_paragraph()
format_bold_lead(ps6_b2, "Objection D : « Je suis complet pour les 6 prochains mois, je n'ai pas besoin de nouveaux chantiers. »\n", 
                 "• Contre-argument 1 : Anticiper pour l'avenir : le référencement naturel (SEO) est un travail de fond. Les articles publiés aujourd'hui génèrent du trafic dans 3 à 6 mois. Commencer maintenant permet de blinder son carnet pour la saison suivante.\n"
                 "• Contre-argument 2 : Choisir de meilleurs chantiers : avoir plus de demandes permet d'être plus sélectif, de choisir les chantiers les plus rentables ou les plus proches géographiquement, et d'augmenter ses marges.")

# ==========================================
# SLIDE 7: Méthodes d'Acquisition Alternatives
# ==========================================
slide7 = prs.slides.add_slide(slide_layout)
apply_background(slide7, COLOR_LIGHT)
add_header(slide7, "Boosters de vente : Canaux d'acquisition et prospection", "STRATÉGIES COMMERCIALES MULTICANAL")

tx_s7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf_s7 = tx_s7.text_frame
tf_s7.word_wrap = True

ps7_intro = tf_s7.paragraphs[0]
ps7_intro.text = "Les consultants commerciaux sont encouragés à utiliser d'autres leviers pour prospecter des experts CVC, Diagnostiqueurs et Bureaux d'Étude :"
ps7_intro.font.bold = True
ps7_intro.font.size = Pt(14)
ps7_intro.font.color.rgb = COLOR_DARK
ps7_intro.space_after = Pt(15)

ps7_b1 = tf_s7.add_paragraph()
format_bold_lead(ps7_b1, "• LinkedIn & Social Selling : ", "Ciblez les dirigeants d'entreprises de climatisation, de plomberie-chauffage ou de cabinets de diagnostic immobilier. Envoyez des messages directs en mettant l'accent sur les problématiques de revente de leads. Exemple : 'Marre de batailler sur les prix avec 4 concurrents sur le même devis ?'.")

ps7_b2 = tf_s7.add_paragraph()
format_bold_lead(ps7_b2, "• Automatisation de la prospection : ", "Utilisez des outils d'extraction et d'emailing pour envoyer des séquences courtes et ciblées aux artisans n'apparaissant pas en première page Google sur leur secteur. Proposez-leur un diagnostic gratuit de leur visibilité locale.")

ps7_b3 = tf_s7.add_paragraph()
format_bold_lead(ps7_b3, "• Campagnes Publicitaires Meta (Facebook / Instagram) : ", "Ciblez les professionnels du BTP en local. [IMPORTANT] Toute création de campagne ou de publicité payante Meta au nom de Gainable.fr doit obligatoirement faire l'objet d'une validation et d'une coordination technique préalable par Exceeddigital.")

# ==========================================
# SLIDE 8: Conclusion & Motivation (Dark Theme)
# ==========================================
slide8 = prs.slides.add_slide(slide_layout)
apply_background(slide8, COLOR_DARK)

tx_s8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0))
tf_s8 = tx_s8.text_frame
tf_s8.word_wrap = True

ps8_title = tf_s8.paragraphs[0]
ps8_title.text = "Votre mission : Devenir le partenaire croissance"
ps8_title.font.name = "Arial"
ps8_title.font.size = Pt(24)
ps8_title.font.bold = True
ps8_title.font.color.rgb = COLOR_WHITE
ps8_title.space_after = Pt(15)

ps8_b1 = tf_s8.add_paragraph()
format_bold_lead(ps8_b1, "✔ Une solution saine : ", "Investissement de visibilité locale durable et éthique.", size=Pt(12.5), font_color=COLOR_WHITE)

ps8_b2 = tf_s8.add_paragraph()
format_bold_lead(ps8_b2, "✔ ROI immédiat : ", "Une adhésion rentabilisée dès le premier chantier signé.", size=Pt(12.5), font_color=COLOR_WHITE)

ps8_b3 = tf_s8.add_paragraph()
format_bold_lead(ps8_b3, "✔ Leviers de prospection libres : ", "Utilisez LinkedIn, phoning, et campagnes Meta coordonnées avec Exceeddigital.", size=Pt(12.5), font_color=COLOR_WHITE)

ps8_b4 = tf_s8.add_paragraph()
format_bold_lead(ps8_b4, "✔ Commissions attractives : ", "Rémunération calculée en temps réel à 17% pour Emy Marty.", size=Pt(12.5), font_color=COLOR_WHITE)

# Right Side: Sales Handshake Photo
if os.path.exists(path_photo_handshake):
    slide8.shapes.add_picture(path_photo_handshake, Inches(7.333), Inches(0), width=Inches(6.0), height=Inches(7.5))

# Save the presentation
output_pptx = os.path.join(project_root, "guide_commercial_gainable.pptx")
prs.save(output_pptx)
print(f"Commercial Guide saved successfully at {output_pptx}")
