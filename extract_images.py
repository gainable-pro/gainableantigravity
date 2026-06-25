from pptx import Presentation
import os

prs = Presentation('gainable_presentation_final.pptx')

# We want to check slides 7 and 8 (which are 0-indexed indices 6 and 7)
for idx in [6, 7]:
    slide = prs.slides[idx]
    print(f"Slide {idx+1} has {len(slide.shapes)} shapes")
    for i, shape in enumerate(slide.shapes):
        if shape.shape_type == 13: # Picture
            image = shape.image
            filename = f"extracted_slide_{idx+1}_img_{i}.{image.ext}"
            with open(filename, 'wb') as f:
                f.write(image.blob)
            print(f"Saved image to {filename} (size: {len(image.blob)} bytes)")
