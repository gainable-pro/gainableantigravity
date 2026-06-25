from pptx import Presentation
import os

prs = Presentation('gainable_presentation_final.pptx')

# We want to check all slides
for idx, slide in enumerate(prs.slides):
    print(f"Slide {idx+1} has {len(slide.shapes)} shapes")
    for i, shape in enumerate(slide.shapes):
        if shape.shape_type == 13: # Picture
            image = shape.image
            filename = f"extracted_all_slide_{idx+1}_img_{i}.{image.ext}"
            with open(filename, 'wb') as f:
                f.write(image.blob)
            print(f"  Saved image: {filename} ({len(image.blob)} bytes)")
