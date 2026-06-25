from pptx import Presentation

prs = Presentation('gainable_presentation_final.pptx')

with open('inspect_out.txt', 'w', encoding='utf-8') as f:
    f.write(f"Total slides: {len(prs.slides)}\n")

    for i, slide in enumerate(prs.slides):
        f.write(f"\n--- Slide {i+1} ---\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        f.write(f"[{shape.name}] {paragraph.text}\n")

