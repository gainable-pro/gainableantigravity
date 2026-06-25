import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Scheme Definitions
COLOR_DARK = RGBColor(31, 45, 61)       # #1F2D3D (Dark slate / navy)
COLOR_LIGHT = RGBColor(248, 250, 252)   # #F8FAFC (Slate-50)
COLOR_GOLD = RGBColor(213, 155, 43)     # #D59B2B (Gainable Gold)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_TEXT_PRIMARY = RGBColor(31, 45, 61)
COLOR_TEXT_SECONDARY = RGBColor(100, 116, 139) # Slate-500
COLOR_LINE = RGBColor(226, 232, 240)    # Slate-200
COLOR_EMERALD = RGBColor(16, 185, 129)  # #10B981

# Paths to directories
project_root = r"C:\Users\gmaro\.gemini\antigravity-ide\scratch\gainable-fr"
assets_dir = r"C:\Users\gmaro\.gemini\antigravity-ide\brain\b1ea4cf3-dabd-4245-96ac-fcaf28419367"
public_dir = os.path.join(project_root, "public")

# Built-in Business & Artisan Photos from 'public/'
path_logo = os.path.join(public_dir, "logo.png")
path_photo_hvac = os.path.join(public_dir, "hero-hvac.png")             # HVAC engineer working
path_photo_diag = os.path.join(public_dir, "diag-inspector.png")        # Real estate inspector
path_photo_vision = os.path.join(public_dir, "espace_pro_vision_1765140841780.png")
path_photo_conclusion = os.path.join(public_dir, "espace_pro_conclusion_1765140856043.png")
path_photo_interior = os.path.join(public_dir, "interior-ac.png")

# Generated Charts
path_chart_seo = os.path.join(project_root, "seo_impact.png")            # SEO growth curve over 3 years
path_chart_visibility = os.path.join(project_root, "seo_curve.png")     # 12-month SEO visibility curve

# Screenshots of Member Account (Air G Energie) & Site (User custom screenshots)
path_scr_homepage = os.path.join(project_root, "user_scr_search.png")
path_scr_profile = os.path.join(project_root, "user_scr_profile.png")
path_scr_leads = os.path.join(assets_dir, "dashboard_leads_1782243573710.png")
path_scr_articles = os.path.join(project_root, "user_scr_articles.png")
path_scr_seo = os.path.join(project_root, "user_scr_seo.png")
path_scr_articles_mobile = os.path.join(project_root, "user_scr_profile_mobile.png")
path_scr_seo_mobile = os.path.join(project_root, "user_scr_profile_mobile.png")

# Extracted slides 7 and 8 screenshots (User custom screenshots)
path_scr_labels = os.path.join(project_root, "user_scr_labels.png")
path_scr_pricing = os.path.join(project_root, "user_scr_pricing.png")

def apply_background(slide, color):
    """Fills the slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, category_text="ACQUISITION CLIENTS"):
    """Adds a standard premium header to light-background slides."""
    # Category prefix (small gold text)
    txBox_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
    tf_cat = txBox_cat.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_GOLD
    
    # Title (large navy text)
    txBox_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf_title = txBox_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Arial"
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_PRIMARY
    
    # Subtle underline border for the header area
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_LINE
    line.line.color.rgb = COLOR_LINE

def format_bold_lead(paragraph, bold_text, normal_text, size=Pt(12.5), font_color=COLOR_TEXT_PRIMARY):
    """Formats a paragraph with bold prefix and regular body text."""
    paragraph.font.name = "Arial"
    paragraph.font.size = size
    
    run_bold = paragraph.add_run()
    run_bold.text = bold_text
    run_bold.font.bold = True
    run_bold.font.color.rgb = font_color
    
    run_normal = paragraph.add_run()
    run_normal.text = normal_text
    run_normal.font.bold = False
    run_normal.font.color.rgb = COLOR_TEXT_SECONDARY

# ==========================================
# SLIDE 1: Title Slide (Premium Half-Image Theme)
# ==========================================
slide_layout = prs.slide_layouts[6] # Blank
slide1 = prs.slides.add_slide(slide_layout)
apply_background(slide1, COLOR_DARK)

# Left Side Panel for Title Text
tx_title = slide1.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(6.5), Inches(4.5))
tf_title = tx_title.text_frame
tf_title.word_wrap = True

# Logo in Title Panel
logo_width = Inches(2.6)
logo_height = Inches(0.78)
slide1.shapes.add_picture(path_logo, Inches(0.8), Inches(1.0), width=logo_width, height=logo_height)

# Main Title
p_t = tf_title.paragraphs[0]
p_t.text = "Gainable.fr"
p_t.font.name = "Arial"
p_t.font.size = Pt(44)
p_t.font.bold = True
p_t.font.color.rgb = COLOR_WHITE
p_t.space_after = Pt(6)

p_t_sub = tf_title.add_paragraph()
p_t_sub.text = "Plateforme de Visibilité & d'Acquisition Locale"
p_t_sub.font.name = "Arial"
p_t_sub.font.size = Pt(22)
p_t_sub.font.bold = True
p_t_sub.font.color.rgb = COLOR_GOLD
p_t_sub.space_after = Pt(20)

p_t_desc = tf_title.add_paragraph()
p_t_desc.text = "Développez votre clientèle CVC, Diagnostics & Bureaux d'Études.\nGénérez des contacts qualifiés en direct, 100% sans commission."
p_t_desc.font.name = "Arial"
p_t_desc.font.size = Pt(14)
p_t_desc.font.color.rgb = COLOR_WHITE
p_t_desc.space_after = Pt(20)

p_t_foot = tf_title.add_paragraph()
p_t_foot.text = "Présentation Commerciale & Guide Pratique"
p_t_foot.font.name = "Arial"
p_t_foot.font.size = Pt(12)
p_t_foot.font.color.rgb = COLOR_TEXT_SECONDARY

# Right Side: Premium HVAC Engineer / Artisan Photo
if os.path.exists(path_photo_hvac):
    slide1.shapes.add_picture(path_photo_hvac, Inches(7.333), Inches(0), width=Inches(6.0), height=Inches(7.5))

# ==========================================
# SLIDE 2: Pourquoi Gainable.fr ? La Transition Numérique (Argumentaire + SEO Impact Curve)
# ==========================================
slide2 = prs.slides.add_slide(slide_layout)
apply_background(slide2, COLOR_LIGHT)
add_header(slide2, "Pourquoi Gainable.fr ? Dominer sa zone à l'ère numérique")

# Left Column: Core Arguments
tx_arg = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.0))
tf_arg = tx_arg.text_frame
tf_arg.word_wrap = True

# Main Intro Paragraph
p_intro = tf_arg.paragraphs[0]
p_intro.text = "Le bouche-à-oreille et les réseaux sociaux sont insuffisants. Face à la concurrence, un référencement Google et IA de premier ordre est devenu primordial."
p_intro.font.name = "Arial"
p_intro.font.size = Pt(14)
p_intro.font.bold = True
p_intro.font.color.rgb = COLOR_DARK
p_intro.space_after = Pt(15)

# Bullet 1: Concurrence & Transition Numérique
p_b1 = tf_arg.add_paragraph()
p_b1.space_after = Pt(12)
format_bold_lead(p_b1, "• Transition Numérique & IA : ", "Les comportements d'achat ont changé. Le bouche-à-oreille ne représente qu'une infime partie du marché. Être présent sur Google et optimisé pour les moteurs de recherche IA est capital.")

# Bullet 2: Objectif Référent Local & Implication
p_b2 = tf_arg.add_paragraph()
p_b2.space_after = Pt(12)
format_bold_lead(p_b2, "• Soyez le référent de votre secteur : ", "Gainable.fr vous fournit l'outil ultime. En publiant régulièrement des articles et en ajoutant des photos de vos chantiers, votre implication et sérieux propulsent votre visibilité locale.")

# Bullet 3: Zéro Commission & Transparence
p_b3 = tf_arg.add_paragraph()
p_b3.space_after = Pt(12)
format_bold_lead(p_b3, "• Zéro commission & Solvabilité validée : ", "Conservez 100% de la valeur de vos chantiers. Lors de l'inscription, une vérification administrative et de solvabilité financière est effectuée via l'API Inessence pour valoriser les vrais professionnels.")

# Right Column: SEO Impact Graph (No CA, cumulative articles vs impressions over 3 years)
if os.path.exists(path_chart_seo):
    slide2.shapes.add_picture(path_chart_seo, Inches(7.0), Inches(2.2), width=Inches(5.5), height=Inches(4.2))

# ==========================================
# SLIDE 3: Le Moteur de Recherche Multi-filtres (Homepage View + Photo Diag)
# ==========================================
slide3 = prs.slides.add_slide(slide_layout)
apply_background(slide3, COLOR_LIGHT)
add_header(slide3, "Le Moteur de Recherche & Contact Client Direct")

# Left Column: Features Explanation
tx_s3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.0))
tf_s3 = tx_s3.text_frame
tf_s3.word_wrap = True

ps3_1 = tf_s3.paragraphs[0]
ps3_1.text = "Un ciblage technique précis pour capter des leads qualifiés"
ps3_1.font.bold = True
ps3_1.font.size = Pt(15)
ps3_1.font.color.rgb = COLOR_DARK
ps3_1.space_after = Pt(15)

ps3_b1 = tf_s3.add_paragraph()
ps3_b1.space_after = Pt(12)
format_bold_lead(ps3_b1, "• Filtres Techniques Précis : ", "Les clients qualifient leur projet par technologie (split, cassette, gainable, VRV...), type d'habitation (maison, appartement...) et par diagnostics (DPE, amiante...), puis filtrent par ville.")

ps3_b2 = tf_s3.add_paragraph()
ps3_b2.space_after = Pt(12)
format_bold_lead(ps3_b2, "• Contact à chaud par téléphone : ", "Permettez aux clients d'accéder directement à vos coordonnées téléphoniques sur votre fiche pour un contact immédiat. Possibilité également de faire des demandes multiples.")

ps3_b3 = tf_s3.add_paragraph()
ps3_b3.space_after = Pt(12)
format_bold_lead(ps3_b3, "• Label Qualité Certifié Expert Gainable : ", "Chaque professionnel sérieux est valorisé par notre badge de certification officiel, renforçant immédiatement la confiance avec le client final.")

# Add a small Business Inspector Photo on the left
if os.path.exists(path_photo_diag):
    slide3.shapes.add_picture(path_photo_diag, Inches(0.8), Inches(4.6), width=Inches(3.2), height=Inches(2.3))

# Right Column: Screenshot Homepage
if os.path.exists(path_scr_homepage):
    slide3.shapes.add_picture(path_scr_homepage, Inches(7.0), Inches(1.8), width=Inches(5.5), height=Inches(4.8))

# ==========================================
# SLIDE 4: Espace Professionnel – Exemple Air G Énergie
# ==========================================
slide4 = prs.slides.add_slide(slide_layout)
apply_background(slide4, COLOR_LIGHT)
add_header(slide4, "L'Espace Pro : Pilotez votre Visibilité & Vos Leads", "ÉTUDE DE CAS : AIR G ENERGIE")

# Left Column: Features Explanation
tx_s4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.0))
tf_s4 = tx_s4.text_frame
tf_s4.word_wrap = True

ps4_1 = tf_s4.paragraphs[0]
ps4_1.text = "Gérez vos informations et réceptionnez vos demandes en direct"
ps4_1.font.bold = True
ps4_1.font.size = Pt(15)
ps4_1.font.color.rgb = COLOR_DARK
ps4_1.space_after = Pt(15)

ps4_b1 = tf_s4.add_paragraph()
ps4_b1.space_after = Pt(12)
format_bold_lead(ps4_b1, "• Fiche Public Référencée : ", "Éditez votre description, vos certifications (RGE...), vos marques partenaires, et vos coordonnées de contact direct.")

ps4_b2 = tf_s4.add_paragraph()
ps4_b2.space_after = Pt(12)
format_bold_lead(ps4_b2, "• Gestionnaire de Leads : ", "Recevez les fiches chantiers qualifiées. Visualisez l'historique complet et contactez directement les clients.")

ps4_b3 = tf_s4.add_paragraph()
ps4_b3.space_after = Pt(12)
format_bold_lead(ps4_b3, "• Facturation & Médias : ", "Ajoutez vos photos de chantiers pour enrichir votre portfolio et gérez vos abonnements Stripe de manière transparente.")

# Add Espace Pro vision photo (Business team meeting / strategy)
if os.path.exists(path_photo_vision):
    slide4.shapes.add_picture(path_photo_vision, Inches(0.8), Inches(4.6), width=Inches(3.2), height=Inches(2.3))

# Right Column: Screenshot Dashboard Profile (Air G Energie) & Mobile mockup overlay
if os.path.exists(path_scr_profile):
    slide4.shapes.add_picture(path_scr_profile, Inches(7.0), Inches(1.8), width=Inches(5.5), height=Inches(4.8))
if os.path.exists(path_scr_articles_mobile):
    # Overlay the mobile mockup screenshot on the right
    slide4.shapes.add_picture(path_scr_articles_mobile, Inches(10.2), Inches(2.2), width=Inches(2.5), height=Inches(4.4))

# ==========================================
# SLIDE 5: Rédaction d'Articles SEO (Le levier de visibilité locale)
# ==========================================
slide5 = prs.slides.add_slide(slide_layout)
apply_background(slide5, COLOR_LIGHT)
add_header(slide5, "La Rédaction d'Articles SEO : Dominez votre Secteur Local", "CONTENU SÉMANTIQUE")

# Left Column: Articles Explanation
tx_s5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
tf_s5 = tx_s5.text_frame
tf_s5.word_wrap = True

ps5_1 = tf_s5.paragraphs[0]
ps5_1.text = "Associez votre expertise à des requêtes géolocalisées"
ps5_1.font.bold = True
ps5_1.font.size = Pt(15)
ps5_1.font.color.rgb = COLOR_DARK
ps5_1.space_after = Pt(15)

ps5_b1 = tf_s5.add_paragraph()
ps5_b1.space_after = Pt(12)
format_bold_lead(ps5_b1, "• Référencement local ultra-ciblé : ", "Rédigez des articles géolocalisés (ex: 'Pose de climatiseur gainable à Miramas') pour apparaître directement devant les internautes de votre secteur.")

ps5_b2 = tf_s5.add_paragraph()
ps5_b2.space_after = Pt(12)
format_bold_lead(ps5_b2, "• Quota Mensuel & Agent IA Expert SEO : ", "Chaque expert dispose d'un quota de 10 articles mensuels. Notre Assistant IA Expert SEO, orienté technique et génie climatique, vous aide à générer vos contenus de haute qualité en quelques minutes.")

ps5_b3 = tf_s5.add_paragraph()
ps5_b3.space_after = Pt(12)
format_bold_lead(ps5_b3, "• Outils Premium d'Édition : ", "Intégrez des images avec balises alt obligatoires pour Google, configurez des FAQ structurées en JSON pour avoir des Rich Snippets, et incorporez des vidéos explicatives.")

# Right Column: Desktop screenshot (User's custom article editor)
if os.path.exists(path_scr_articles):
    slide5.shapes.add_picture(path_scr_articles, Inches(7.0), Inches(1.8), width=Inches(5.5), height=Inches(4.8))

# ==========================================
# SLIDE 6: Suivi de Performance & Concurrence (SEO GSC)
# ==========================================
slide6 = prs.slides.add_slide(slide_layout)
apply_background(slide6, COLOR_LIGHT)
add_header(slide6, "Suivi de Performance & Comparateur de Concurrence", "GOOGLE SEARCH CONSOLE & BENCHMARKS")

# Left Column: GSC Metrics & Explanations
tx_s6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.0))
tf_s6 = tx_s6.text_frame
tf_s6.word_wrap = True

ps6_1 = tf_s6.paragraphs[0]
ps6_1.text = "L'analyse scientifique de vos performances SEO"
ps6_1.font.bold = True
ps6_1.font.size = Pt(15)
ps6_1.font.color.rgb = COLOR_DARK
ps6_1.space_after = Pt(15)

ps6_b1 = tf_s6.add_paragraph()
ps6_b1.space_after = Pt(12)
format_bold_lead(ps6_b1, "• Graphique GSC natif : ", "Visualisez vos impressions, vos clics totaux, votre taux de clic (CTR moyen) et la position moyenne de vos mots-clés directement dans votre Espace Pro.")

ps6_b2 = tf_s6.add_paragraph()
ps6_b2.space_after = Pt(12)
format_bold_lead(ps6_b2, "• Diagnostic & Score de votre Site : ", "Obtenez un score sur 100 mesurant la conformité technique de votre site web personnel (sécurité HSTS, Content Security Policy, en-têtes...).")

ps6_b3 = tf_s6.add_paragraph()
ps6_b3.space_after = Pt(12)
format_bold_lead(ps6_b3, "• Comparateur face aux leaders nationaux : ", "Suivez en temps réel comment votre page Gainable.fr se positionne face aux géants (Maclem, IZI EDF, Engie...) et profitez de l'autorité globale de la plateforme.")

# Right Column: Desktop screenshot & Plotted SEO visibility curve
if os.path.exists(path_scr_seo):
    slide6.shapes.add_picture(path_scr_seo, Inches(6.8), Inches(1.8), width=Inches(5.5), height=Inches(4.8))
if os.path.exists(path_chart_visibility):
    # Position the 12-month SEO curve chart overlapping at the bottom-right
    slide6.shapes.add_picture(path_chart_visibility, Inches(9.5), Inches(3.2), width=Inches(3.2), height=Inches(2.5))

# ==========================================
# SLIDE 7: Le Label Gainable.fr (Gage de Confiance & Qualité - User Slide 7)
# ==========================================
slide7 = prs.slides.add_slide(slide_layout)
apply_background(slide7, COLOR_LIGHT)
add_header(slide7, "Le Label Gainable.fr : Gage de Confiance & Transparence", "LABELLISATION")

# Left Column: Label Explanation
tx_s7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.0))
tf_s7 = tx_s7.text_frame
tf_s7.word_wrap = True

ps7_1 = tf_s7.paragraphs[0]
ps7_1.text = "Valorisez votre savoir-faire auprès de vos futurs clients"
ps7_1.font.bold = True
ps7_1.font.size = Pt(15)
ps7_1.font.color.rgb = COLOR_DARK
ps7_1.space_after = Pt(15)

ps7_b1 = tf_s7.add_paragraph()
ps7_b1.space_after = Pt(12)
format_bold_lead(ps7_b1, "• Transparence & Confiance : ", "Gainable.fr valorise les professionnels sérieux et garantit aux particuliers une recherche simple, rapide et fiable.")

ps7_b2 = tf_s7.add_paragraph()
ps7_b2.space_after = Pt(12)
format_bold_lead(ps7_b2, "• Badge \"Expert Vérifié\" visible : ", "Ce label, gage de qualité, s'affiche directement sur votre profil public pour maximiser la conversion de vos prospects en clients.")

ps7_b3 = tf_s7.add_paragraph()
ps7_b3.space_after = Pt(12)
format_bold_lead(ps7_b3, "• Critères de Labellisation stricts : ", "Vérification des informations administratives, contrôle des marques et certifications, validation de la zone géographique et analyse de solvabilité via API Inessence.")

# Right Column: Screenshot Labels Page
if os.path.exists(path_scr_labels):
    slide7.shapes.add_picture(path_scr_labels, Inches(7.0), Inches(1.8), width=Inches(5.5), height=Inches(4.8))

# ==========================================
# SLIDE 8: Tarifs transparents & Formules d'Adhésion (User Slide 8)
# ==========================================
slide8 = prs.slides.add_slide(slide_layout)
apply_background(slide8, COLOR_LIGHT)
add_header(slide8, "Tarifs Transparents & Formules d'Adhésion", "FORMULES ET OFFRES")

# Left Column: Pricing Explanation
tx_s8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.0))
tf_s8 = tx_s8.text_frame
tf_s8.word_wrap = True

ps8_1 = tf_s8.paragraphs[0]
ps8_1.text = "Des offres adaptées à votre profil de professionnel"
ps8_1.font.bold = True
ps8_1.font.size = Pt(15)
ps8_1.font.color.rgb = COLOR_DARK
ps8_1.space_after = Pt(15)

ps8_b1 = tf_s8.add_paragraph()
ps8_b1.space_after = Pt(12)
format_bold_lead(ps8_b1, "• Société Expert CVC : ", "850 € HT par an (Tarif de lancement au lieu de 1200 € HT). Conçu spécifiquement pour les installateurs de climatisation, clim réversible et PAC.")

ps8_b2 = tf_s8.add_paragraph()
ps8_b2.space_after = Pt(12)
format_bold_lead(ps8_b2, "• Diagnostiqueur Immobilier : ", "750 € HT par an (Tarif de lancement au lieu de 1100 € HT). Pour les experts en diagnostics DPE, amiante, etc.")

ps8_b3 = tf_s8.add_paragraph()
ps8_b3.space_after = Pt(12)
format_bold_lead(ps8_b3, "• Bureau d'Étude Thermique : ", "Gratuit pour les experts en thermique (ciblage RE2020, codes APE BE dédiés).")

ps8_b4 = tf_s8.add_paragraph()
ps8_b4.space_after = Pt(12)
format_bold_lead(ps8_b4, "• Toutes fonctionnalités incluses : ", "0% de commission sur vos contacts, leads et demandes illimités, quota de 10 articles SEO rédigés par mois, assistant IA, et badge de labellisation.")

# Right Column: Screenshot Pricing Page
if os.path.exists(path_scr_pricing):
    slide8.shapes.add_picture(path_scr_pricing, Inches(7.0), Inches(1.8), width=Inches(5.5), height=Inches(4.8))

# ==========================================
# SLIDE 9: Conclusion & Call to Action (Dark Theme + Photo Artisan - User Slide 9)
# ==========================================
slide9 = prs.slides.add_slide(slide_layout)
apply_background(slide9, COLOR_DARK)

# Call to Action Title
tx_c = slide9.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6.5), Inches(5.5))
tf_c = tx_c.text_frame
tf_c.word_wrap = True

pc_t = tf_c.paragraphs[0]
pc_t.text = "Prenez le contrôle de votre Acquisition Locale"
pc_t.font.name = "Arial"
pc_t.font.size = Pt(26)
pc_t.font.bold = True
pc_t.font.color.rgb = COLOR_WHITE
pc_t.space_after = Pt(8)

pc_t2 = tf_c.add_paragraph()
pc_t2.text = "Boostez vos demandes de climatisation, diagnostics et études thermiques sans intermédiaires."
pc_t2.font.name = "Arial"
pc_t2.font.size = Pt(13.5)
pc_t2.font.color.rgb = COLOR_GOLD
pc_t2.space_after = Pt(15)

# Bullet points as requested by the user's Slide 9 content
pc_b1 = tf_c.add_paragraph()
pc_b1.text = "✔ Abonnement fixe transparent & 0% commission"
pc_b1.font.size = Pt(13)
pc_b1.font.color.rgb = COLOR_WHITE
pc_b1.space_after = Pt(6)

pc_b2 = tf_c.add_paragraph()
pc_b2.text = "✔ Valoriser les vrais professionnels"
pc_b2.font.size = Pt(13)
pc_b2.font.color.rgb = COLOR_WHITE
pc_b2.space_after = Pt(6)

pc_b3 = tf_c.add_paragraph()
pc_b3.text = "✔ Apporter des contacts réellement qualifiés"
pc_b3.font.size = Pt(13)
pc_b3.font.color.rgb = COLOR_WHITE
pc_b3.space_after = Pt(6)

pc_b4 = tf_c.add_paragraph()
pc_b4.text = "✔ Rétablir la confiance entre artisans et clients"
pc_b4.font.size = Pt(13)
pc_b4.font.color.rgb = COLOR_WHITE
pc_b4.space_after = Pt(6)

pc_b5 = tf_c.add_paragraph()
pc_b5.text = "✔ Donner aux TPE/PME les outils pour rivaliser"
pc_b5.font.size = Pt(13)
pc_b5.font.color.rgb = COLOR_WHITE
pc_b5.space_after = Pt(6)

pc_b6 = tf_c.add_paragraph()
pc_b6.text = "✔ Plus de 58 000 pages et articles référencés sur Google"
pc_b6.font.size = Pt(13)
pc_b6.font.color.rgb = COLOR_WHITE
pc_b6.space_after = Pt(6)

pc_b7 = tf_c.add_paragraph()
pc_b7.text = "✔ Suivi en temps réel des leads et impressions"
pc_b7.font.size = Pt(13)
pc_b7.font.color.rgb = COLOR_WHITE
pc_b7.space_after = Pt(15)

# Contact info (Gainable sites & email)
pb_contact = tf_c.add_paragraph()
pb_contact.text = "Rejoignez le réseau : www.gainable.fr | www.gainable.be | www.gainable.ch\nContact commercial : contact@gainable.ch"
pb_contact.font.name = "Arial"
pb_contact.font.size = Pt(12.5)
pb_contact.font.bold = True
pb_contact.font.color.rgb = COLOR_GOLD

# Right Side: Professional Conclusion / Team / Artisan Photo
if os.path.exists(path_photo_conclusion):
    slide9.shapes.add_picture(path_photo_conclusion, Inches(7.333), Inches(0), width=Inches(6.0), height=Inches(7.5))

# Save presentation
try:
    output_pptx = os.path.join(project_root, "gainable_presentation_final.pptx")
    prs.save(output_pptx)
    print(f"Presentation saved successfully at {output_pptx}")
except PermissionError:
    output_pptx = os.path.join(project_root, "gainable_presentation_final_updated.pptx")
    prs.save(output_pptx)
    print(f"Presentation saved successfully (with lock fallback) at {output_pptx}")
